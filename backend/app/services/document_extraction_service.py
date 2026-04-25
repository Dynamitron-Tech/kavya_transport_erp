"""
Kavya Transports — Document Extraction Service
Uses Claude Sonnet 4.6 vision to extract structured data
from uploaded transport documents.
"""

import anthropic
import base64
import json
import os
import re
import io
import logging
import httpx
from datetime import datetime
from app.core.config import settings

logger = logging.getLogger(__name__)


# ─────────────────────────────────────────────────────────────
# EXTRACTION PROMPTS — one per uploadable document type
# ─────────────────────────────────────────────────────────────

EXTRACTION_PROMPTS: dict[str, str] = {

    "rc": """
You are a document data extraction assistant. Extract specific fields from an Indian vehicle Registration Certificate (RC) smart card.

This is typically a Tamil Nadu / Indian government-issued RC smart card with a chip. Read ALL text carefully.

FIELDS TO EXTRACT — use the EXACT label names listed:

1. "registration_number"
   - Label on card: "Regn. Number" or "Reg. No." or at top of card
   - Example: TN72BC7214 (remove all spaces)

2. "owner_name"
   - Label on card: "Owner Name"
   - The value is the FULL NAME printed on the next line BELOW the "Owner Name" label
   - Example: "NALAN SHUNMUGARAJ K"
   - ONLY letters and spaces — never numbers or codes
   - Do NOT return the label word "Owner" itself — return the actual person/company name

3. "vehicle_class"
   - Label on card: "Vehicle Class" or "Class of Vehicle"
   - Example: "LMV", "HGV", "TRANS"

4. "fuel_type"
   - Label on card: "Fuel" (usually on left side of card)
   - Example: "DIESEL" → normalise to "Diesel"
   - Normalise to one of: Diesel / Petrol / CNG / Electric / LPG

5. "engine_number"
   - Label on card: "Engine/Motor Number" (exact label used on Tamil Nadu RC)
   - The value is the alphanumeric code printed on the NEXT LINE below "Engine/Motor Number"
   - Example: "1ND1440167"
   - MUST be alphanumeric (letters + digits). NEVER a person's name.

6. "chassis_number"
   - Label on card: "Chassis Number" (exact label used on Tamil Nadu RC)
   - The value is the alphanumeric VIN code printed on the NEXT LINE below "Chassis Number"
   - Example: "MBJB49BT90001213701215"
   - MUST be a long alphanumeric code (17+ characters). NEVER a person's name.
   - CRITICAL: The chassis number appears BEFORE "Engine/Motor Number" on the card. Do NOT confuse with owner name.

7. "issue_date"
   - Label on card: "Date of Regn." or "Regn. Date"
   - Dates on card may use dashes (04-01-2016) — convert to DD/MM/YYYY format → "04/01/2016"

8. "validity_date"
   - Label on card: "Regn. Validity" or "Valid Upto"
   - Convert to DD/MM/YYYY format

CRITICAL RULES:
- Read the label first, then extract the value on the NEXT LINE after that label
- chassis_number and engine_number are ALWAYS long alphanumeric codes, NEVER names
- owner_name is ALWAYS a person or company name (letters only), NEVER a code
- Dates with dashes (DD-MM-YYYY) must be converted to slashes (DD/MM/YYYY)
- Return ONLY valid JSON with no extra text or markdown

Return exactly:
{"registration_number": "...", "owner_name": "...", "vehicle_class": "...", "fuel_type": "...", "engine_number": "...", "chassis_number": "...", "issue_date": "DD/MM/YYYY", "validity_date": "DD/MM/YYYY"}

Set any not-found field to null.
""",

    "insurance": """
You are a precise document data extraction assistant specialised
in Indian transport documents. Extract data from this
Vehicle Insurance Certificate image.

FIELDS TO EXTRACT:
1. policy_number  — Insurance policy number
2. insurer_name   — Name of the insurance company
3. issue_date     — Date the policy was issued / started
4. expiry_date    — Date the policy expires
5. vehicle_number — Vehicle registration number this policy covers

EXTRACTION RULES:
- policy_number: appears near labels like "Policy No",
  "Certificate No", "Policy Number".
- insurer_name: the insurance company full official name.
- issue_date / expiry_date: both must be in DD/MM/YYYY format.
  "From" date = issue_date, "To" date = expiry_date.
- vehicle_number: the vehicle registration number.

OUTPUT: Return ONLY a valid JSON object. No markdown.

{
  "policy_number": "...",
  "insurer_name": "...",
  "issue_date": "DD/MM/YYYY",
  "expiry_date": "DD/MM/YYYY",
  "vehicle_number": "..."
}

Set any field you cannot find to null. Do not guess.
""",

    "fitness": """
You are a precise document data extraction assistant specialised
in Indian transport documents. Extract data from this
Vehicle Fitness Certificate image.

FIELDS TO EXTRACT:
1. certificate_number — Fitness certificate number
2. vehicle_number     — Vehicle registration number
3. issue_date         — Date of issue
4. expiry_date        — Date of expiry

EXTRACTION RULES:
- certificate_number: appears near labels like "Cert No",
  "Fitness Cert No", "FC No".
- All dates: normalise to DD/MM/YYYY format.

OUTPUT: Return ONLY a valid JSON object. No markdown.

{
  "certificate_number": "...",
  "vehicle_number": "...",
  "issue_date": "DD/MM/YYYY",
  "expiry_date": "DD/MM/YYYY"
}

Set any field you cannot find to null. Do not guess.
""",

    "puc": """
You are a precise document data extraction assistant specialised
in Indian transport documents. Extract data from this
Pollution Under Control (PUC) Certificate image.

FIELDS TO EXTRACT:
1. certificate_number — PUC certificate number
2. vehicle_number     — Vehicle registration number
3. test_date          — Date the emission test was conducted
4. valid_until        — Certificate validity / expiry date
5. reading_values     — Emission reading values from the test

EXTRACTION RULES:
- certificate_number: appears near labels like "Cert No", "PUC No".
- test_date / valid_until: normalise to DD/MM/YYYY format.
- reading_values: extract ALL emission readings as a nested object.
  Common readings:
    CO (Carbon Monoxide) → key "co_percent"
    HC (Hydrocarbons)    → key "hc_ppm"
    CO2                  → key "co2_percent"
    Smoke opacity        → key "smoke_opacity"
  If no readings found, return null for this field.

OUTPUT: Return ONLY a valid JSON object. No markdown.

{
  "certificate_number": "...",
  "vehicle_number": "...",
  "test_date": "DD/MM/YYYY",
  "valid_until": "DD/MM/YYYY",
  "reading_values": {"co_percent": "0.15", "hc_ppm": "85"}
}

Set any top-level field you cannot find to null. Do not guess.
""",

    "permit": """
You are a precise document data extraction assistant specialised
in Indian transport documents. Extract data from this
Vehicle Permit document image.

FIELDS TO EXTRACT:
1. permit_number  — Permit number / authorisation number
2. vehicle_number — Vehicle registration number
3. route_area     — Route or area this permit covers
4. issue_date     — Date of issue
5. expiry_date    — Date of expiry

EXTRACTION RULES:
- permit_number: appears near labels like "Permit No", "P.No".
- route_area: the permitted route or area (full text, max 200 chars).
- All dates: normalise to DD/MM/YYYY format.

OUTPUT: Return ONLY a valid JSON object. No markdown.

{
  "permit_number": "...",
  "vehicle_number": "...",
  "route_area": "...",
  "issue_date": "DD/MM/YYYY",
  "expiry_date": "DD/MM/YYYY"
}

Set any field you cannot find to null. Do not guess.
""",

    "tax_receipt": """
You are a precise document data extraction assistant specialised
in Indian transport documents. Extract data from this
Road Tax Receipt image.

FIELDS TO EXTRACT:
1. receipt_number  — Receipt / challan / transaction number
2. vehicle_number  — Vehicle registration number
3. tax_amount      — Amount of road tax paid
4. valid_until     — Tax validity / next due date

EXTRACTION RULES:
- receipt_number: appears near labels like "Receipt No",
  "Transaction No", "Challan No", "MR No".
- tax_amount: include currency symbol if present.
- valid_until: normalise to DD/MM/YYYY format.

OUTPUT: Return ONLY a valid JSON object. No markdown.

{
  "receipt_number": "...",
  "vehicle_number": "...",
  "tax_amount": "₹5,200",
  "valid_until": "DD/MM/YYYY"
}

Set any field you cannot find to null. Do not guess.
""",

        "driving_license": """
You are a document data extraction assistant. Extract the following fields from a Driving Licence (DL) document.

Fields to extract:
1. License Number
2. Issue Date
3. Validity Date

Rules:
- "license_number": The DL number may appear in two ways:
    a) At the TOP of the card with NO label — extract the first prominent alphanumeric value at the top.
    b) With a label such as "License No", "DL No", "Driving Licence Number", "No.", "DL#", or similar — extract the value next to that label.
  If both exist, prefer the labelled one. Remove all spaces from the number (e.g. "TN72 20240005499" → "TN7220240005499").
- "issue_date": look ONLY for labels explicitly saying "Issue Date" or "Date of Issue". This is the date the licence was issued/printed (typically a recent date). IMPORTANT: Do NOT use "Date of Birth", "DOB", or "D.O.B" as the issue date — those are different fields. Normalise to DD/MM/YYYY.
- "expiry_date": look for labels "Valid Upto", "Validity", "Validity(NT)", "Validity(TR)". If multiple validity dates exist, use the LATEST one. Normalise to DD/MM/YYYY.
- Dates may appear in any format (DD/MM/YYYY, DD-MM-YYYY, DD MMM YYYY, etc.) — always normalise to DD/MM/YYYY.

Return ONLY a JSON object with keys: "license_number", "issue_date", "expiry_date"
Set any not-found field to null. No explanation, no extra text.
""",

    "aadhaar": """
You are a precise document data extraction assistant specialised
in Indian government identity documents. Extract data from this
Aadhaar Card image.

FIELDS TO EXTRACT:
1. aadhaar_number — 12-digit Aadhaar identification number
2. date_of_birth  — Date of birth
3. gender         — Gender of the holder

EXTRACTION RULES:
- aadhaar_number: always 12 digits, format "1234 5678 9012".
  Extract ONLY digits as 12-char string: "123456789012".
  If first 8 digits masked: return "XXXXXXXX" + last4.
- date_of_birth: normalise to DD/MM/YYYY. If only year: return "1990".
- gender: normalise to "Male", "Female", or "Transgender".

PRIVACY RULE: Do NOT describe the photo, QR code, or biometrics.
Do NOT extract the name or address — only the fields listed above.

OUTPUT: Return ONLY a valid JSON object. No markdown.

{
  "aadhaar_number": "123456789012",
  "date_of_birth": "15/06/1990",
  "gender": "Male"
}

Set any field you cannot find to null. Do not guess.
""",

    "driver_badge": """
You are a precise document data extraction assistant specialised
in Indian transport documents. Extract data from this
Driver Badge or Driver ID Card image.

FIELDS TO EXTRACT:
1. badge_number             — Badge number / ID number
2. driver_name              — Full name of the driver
3. issuing_authority        — Company or authority that issued this badge
4. issue_date               — Date the badge was issued
5. expiry_date              — Expiry / renewal date of the badge
6. vehicle_type_authorized  — Type of vehicle this driver is authorised to operate

EXTRACTION RULES:
- badge_number: appears near "Badge No", "ID No", "Driver ID".
- issue_date / expiry_date: normalise to DD/MM/YYYY format.
- vehicle_type_authorized: may say "Truck", "Tanker", "HGV" etc.

OUTPUT: Return ONLY a valid JSON object. No markdown.

{
  "badge_number": "...",
  "driver_name": "...",
  "issuing_authority": "...",
  "issue_date": "DD/MM/YYYY",
  "expiry_date": "DD/MM/YYYY",
  "vehicle_type_authorized": "..."
}

Set any field you cannot find to null. Do not guess.
""",

    "medical_fitness": """
You are a precise document data extraction assistant specialised
in Indian transport documents. Extract data from this
Medical Fitness Certificate for a commercial vehicle driver.

FIELDS TO EXTRACT:
1. certificate_number — Certificate / registration number
2. driver_name        — Full name of the driver examined
3. date_of_birth      — Date of birth of the driver
4. examination_date   — Date the medical examination was done
5. valid_until        — Certificate validity / expiry date
6. doctor_name        — Name of the issuing doctor
7. hospital_name      — Name of the hospital or clinic
8. fitness_status     — Whether the driver is Fit or Unfit

EXTRACTION RULES:
- All dates: normalise to DD/MM/YYYY format.
- doctor_name: include "Dr." prefix if printed.
- fitness_status: normalise to exactly one of:
  "Fit", "Unfit", or "Conditionally Fit".

OUTPUT: Return ONLY a valid JSON object. No markdown.

{
  "certificate_number": "...",
  "driver_name": "...",
  "date_of_birth": "DD/MM/YYYY",
  "examination_date": "DD/MM/YYYY",
  "valid_until": "DD/MM/YYYY",
  "doctor_name": "Dr. ...",
  "hospital_name": "...",
  "fitness_status": "Fit"
}

Set any field you cannot find to null. Do not guess.
""",

    "pan_card": """
You are a precise document data extraction assistant specialised
in Indian government documents. Extract data from this
PAN Card (Permanent Account Number) image.

FIELDS TO EXTRACT:
1. pan_number      — 10-character PAN alphanumeric code
2. full_name       — Full name of the PAN holder
3. date_of_birth   — Date of birth or date of incorporation
4. pan_holder_type — Type of PAN holder (inferred from 4th character)

EXTRACTION RULES:
- pan_number: always 10 chars, format 5 letters + 4 digits + 1 letter.
  Extract as uppercase with no spaces (e.g., ABCDE1234F).
- full_name: English name exactly as printed.
- date_of_birth: normalise to DD/MM/YYYY format.
- pan_holder_type: infer from 4th character of PAN:
  P → "Individual", C → "Company", H → "Hindu Undivided Family",
  F → "Firm / Partnership", A → "Association of Persons",
  B → "Body of Individuals", G → "Government",
  J → "Artificial Juridical Person", L → "Local Authority",
  T → "Trust"

OUTPUT: Return ONLY a valid JSON object. No markdown.

{
  "pan_number": "ABCDE1234F",
  "full_name": "Kumar Selvam",
  "date_of_birth": "15/06/1990",
  "pan_holder_type": "Individual"
}

Set any field you cannot find to null. Do not guess.
""",

    "gst_certificate": """
You are a precise document data extraction assistant specialised
in Indian government business documents. Extract data from this
GST Registration Certificate image.

FIELDS TO EXTRACT:
1. gstin             — GST Identification Number (15 characters)
2. legal_name        — Legal name of the registered business
3. trade_name        — Trade name (if explicitly labelled)
4. registration_date — Date of GST registration
5. business_type     — Type of business entity
6. principal_address — Principal place of business address
7. gst_status        — Current status of the GST registration

EXTRACTION RULES:
- gstin: always 15 characters, uppercase, no spaces.
  Format: 2-digit state code + 10-char PAN + entity no + check chars.
- trade_name: return null if not explicitly labelled "Trade Name".
- registration_date: normalise to DD/MM/YYYY format.
- business_type: normalise to one of:
  "Proprietorship", "Partnership", "LLP", "Private Limited",
  "Public Limited", "Trust", "Society", "Government", "Other"
- gst_status: normalise to one of:
  "Active", "Inactive", "Cancelled", "Suspended"

OUTPUT: Return ONLY a valid JSON object. No markdown.

{
  "gstin": "33AABCT1332L1ZD",
  "legal_name": "...",
  "trade_name": null,
  "registration_date": "DD/MM/YYYY",
  "business_type": "Private Limited",
  "principal_address": "...",
  "gst_status": "Active"
}

Set any field you cannot find to null. Do not guess.
""",

    "passbook": """
You are an intelligent document parser specialised in Indian bank passbooks.
Extract structured banking details from the passbook image.

FIELDS TO EXTRACT:
1. account_holder_name  — Full name of the account holder (UPPERCASE, no titles)
2. account_number       — Digits only, 9–18 digits, remove spaces/symbols
3. bank_name            — Full official bank name (e.g. STATE BANK OF INDIA)
4. branch_name          — Branch location/name only
5. ifsc_code            — Must match pattern: 4 letters + 0 + 6 alphanumeric
6. micr_code            — 9-digit MICR number (optional)
7. bank_address         — Full branch address (optional, 1-2 lines)

EXTRACTION RULES:
- account_holder_name: look for "Name", "A/c Name", "Account Holder". Output UPPERCASE.
- account_number: look for "Account No", "A/c No", "Account Number". Digits only.
- bank_name: top/header of document. Normalise to full official name.
- branch_name: look for "Branch", "Branch Name". Location only.
- ifsc_code: look for "IFSC". Strict format [A-Z]{4}0[A-Z0-9]{6}.
- micr_code: look for "MICR". 9-digit number.
- bank_address: only if clearly labelled.

VALIDATION:
- IFSC must strictly match the pattern or return null.
- Account number must be numeric only.
- Name must not contain numbers.
- If validation fails, return best guess and set confidence to "low".

OUTPUT: Return ONLY a valid JSON object. No markdown. No explanations.

{
  "account_holder_name": "KUMAR SELVAM",
  "account_number": "123456789012",
  "bank_name": "STATE BANK OF INDIA",
  "branch_name": "ANNA NAGAR",
  "ifsc_code": "SBIN0001234",
  "micr_code": "600002123",
  "bank_address": "12 Anna Nagar, Chennai 600040",
  "confidence": "high"
}

Set any field you cannot find to null. Do not guess.
""",
}



# ─────────────────────────────────────────────────────────────
# SYSTEM-GENERATED document types — no extraction performed
# ─────────────────────────────────────────────────────────────

SYSTEM_GENERATED_TYPES: set[str] = {
    "invoice",
    "eway_bill",
    "lr_copy",
    "contract",
    "pod",
    "other",
}

DOCUMENT_TYPE_ALIASES: dict[str, str] = {
    "license": "driving_license",
    "bank_passbook": "passbook",
    "bank_account": "passbook",
}


# ─────────────────────────────────────────────────────────────
# ENTITY → DOCUMENT REQUIREMENTS CONFIG
# ─────────────────────────────────────────────────────────────

ENTITY_DOCUMENTS: dict = {
    "vehicle": {
        "required": [
            {
                "type": "rc",
                "label": "Registration Certificate (RC)",
                "description": "Vehicle registration issued by RTO",
                "icon": "file-text",
                "auto_fills": ["registration_number", "fuel_type", "vehicle_class"],
            },
            {
                "type": "insurance",
                "label": "Insurance Certificate",
                "description": "Comprehensive or third-party insurance",
                "icon": "shield",
                "auto_fills": [],
            },
            {
                "type": "fitness",
                "label": "Fitness Certificate",
                "description": "Vehicle roadworthiness certificate from RTO",
                "icon": "check-circle",
                "auto_fills": [],
            },
            {
                "type": "puc",
                "label": "PUC Certificate",
                "description": "Pollution under control emission certificate",
                "icon": "wind",
                "auto_fills": [],
            },
        ],
        "optional": [
            {
                "type": "permit",
                "label": "Permit",
                "description": "Commercial vehicle permit — state / all-India",
                "icon": "map",
                "auto_fills": [],
            },
            {
                "type": "tax_receipt",
                "label": "Road Tax Receipt",
                "description": "Road tax payment receipt",
                "icon": "receipt",
                "auto_fills": [],
            },
        ],
    },
    "driver": {
        "required": [
            {
                "type": "driving_license",
                "label": "Driving License",
                "description": "Valid commercial driving license",
                "icon": "credit-card",
                "auto_fills": ["license_number", "date_of_birth", "license_classes"],
            },
            {
                "type": "aadhaar",
                "label": "Aadhaar Card",
                "description": "Government-issued identity proof",
                "icon": "user",
                "auto_fills": ["full_name", "date_of_birth", "address"],
            },
        ],
        "optional": [],
    },
    "employee": {
        "required": [
            {
                "type": "aadhaar",
                "label": "Aadhaar Card",
                "description": "Government-issued identity proof",
                "icon": "user",
                "auto_fills": ["full_name", "date_of_birth", "address"],
            },
        ],
        "optional": [
            {
                "type": "pan_card",
                "label": "PAN Card",
                "description": "Required for salary TDS deduction",
                "icon": "credit-card",
                "auto_fills": [],
            },
        ],
    },
    "client": {
        "required": [
            {
                "type": "gst_certificate",
                "label": "GST Registration Certificate",
                "description": "GST registration — required for billing and e-way bills",
                "icon": "file-badge",
                "auto_fills": ["gstin", "legal_name", "principal_address"],
            },
        ],
        "optional": [
            {
                "type": "pan_card",
                "label": "PAN Card",
                "description": "PAN for TDS compliance",
                "icon": "credit-card",
                "auto_fills": [],
            },
        ],
    },
}

# ─────────────────────────────────────────────────────────────
# EXPIRY FIELD NAMES per document type
# ─────────────────────────────────────────────────────────────

EXPIRY_FIELDS: dict[str, str] = {
    "rc": "validity_date",
    "insurance": "expiry_date",
    "fitness": "expiry_date",
    "puc": "valid_until",
    "permit": "expiry_date",
    "tax_receipt": "valid_until",
    "driving_license": "expiry_date",
    "driver_badge": "expiry_date",
    "medical_fitness": "valid_until",
}

ISSUE_DATE_FIELDS: dict[str, str] = {
    "rc": "issue_date",
    "insurance": "issue_date",
    "fitness": "issue_date",
    "puc": "test_date",
    "permit": "issue_date",
    "driving_license": "issue_date",
    "driver_badge": "issue_date",
    "medical_fitness": "examination_date",
    "gst_certificate": "registration_date",
}

# Mapping from API lowercase type → DB DocumentType enum value
DOC_TYPE_TO_ENUM: dict[str, str] = {
    "rc": "RC",
    "insurance": "INSURANCE",
    "fitness": "FITNESS",
    "puc": "PUC",
    "permit": "PERMIT",
    "tax_receipt": "TAX_RECEIPT",
    "driving_license": "LICENSE",
    "aadhaar": "AADHAAR",
    "driver_badge": "DRIVER_BADGE",
    "medical_fitness": "MEDICAL_FITNESS",
    "pan_card": "PAN_CARD",
    "gst_certificate": "GST_CERTIFICATE",
    "invoice": "INVOICE",
    "eway_bill": "EWAY_BILL",
    "lr_copy": "LR_COPY",
    "contract": "CONTRACT",
    "pod": "POD",
    "other": "OTHER",
}


# ─────────────────────────────────────────────────────────────
# Donut DocVQA field prompts
# ─────────────────────────────────────────────────────────────

DONUT_FIELD_QUESTIONS: dict[str, dict[str, str]] = {
    "driving_license": {
        "license_number": "What is the driving license number?",
        "issue_date": "What is the issue date of this driving license?",
        "expiry_date": "What is the validity expiry date of this driving license?",
    },
    "rc": {
        "registration_number": "What is the vehicle registration number?",
        "owner_name": "What is the owner name?",
        "vehicle_class": "What is the vehicle class?",
        "fuel_type": "What is the fuel type?",
        "engine_number": "What is the engine number?",
        "chassis_number": "What is the chassis number?",
        "issue_date": "What is the registration issue date?",
        "validity_date": "What is the registration validity date?",
    },
    "insurance": {
        "policy_number": "What is the insurance policy number?",
        "insurer_name": "What is the insurer name?",
        "issue_date": "What is the policy issue date?",
        "expiry_date": "What is the policy expiry date?",
        "vehicle_number": "What is the vehicle registration number?",
    },
    "fitness": {
        "certificate_number": "What is the fitness certificate number?",
        "vehicle_number": "What is the vehicle registration number?",
        "issue_date": "What is the issue date?",
        "expiry_date": "What is the expiry date?",
    },
    "puc": {
        "certificate_number": "What is the PUC certificate number?",
        "vehicle_number": "What is the vehicle registration number?",
        "test_date": "What is the test date?",
        "valid_until": "What is the validity end date?",
    },
    "permit": {
        "permit_number": "What is the permit number?",
        "vehicle_number": "What is the vehicle registration number?",
        "route_area": "What is the route or area of operation?",
        "issue_date": "What is the permit issue date?",
        "expiry_date": "What is the permit expiry date?",
    },
    "tax_receipt": {
        "receipt_number": "What is the receipt or challan number?",
        "vehicle_number": "What is the vehicle registration number?",
        "tax_amount": "What is the tax amount paid?",
        "valid_until": "What is the tax validity date?",
    },
    "aadhaar": {
        "date_of_birth": "What is the date of birth on this Aadhaar card?",
        "gender": "What is the gender on this Aadhaar card?",
        "aadhaar_number": "What is the 12-digit Aadhaar number?",
    },
    "passbook": {
        "account_holder_name": "What is the account holder name on this bank passbook?",
        "account_number": "What is the account number on this bank passbook?",
        "bank_name": "What is the bank name on this passbook?",
        "branch_name": "What is the branch name on this passbook?",
        "ifsc_code": "What is the IFSC code on this passbook?",
        "micr_code": "What is the MICR code on this passbook?",
    },
}

PRESENTATION_MIME_TYPES: set[str] = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.ms-powerpoint",
}


# ─────────────────────────────────────────────────────────────
# DOCUMENT EXTRACTION SERVICE
# ─────────────────────────────────────────────────────────────

class DocumentExtractionService:
    """
    Document extraction service with multiple backends:
    1. Hugging Face Inference API (Donut DocVQA) — preferred, API-based
    2. Claude Sonnet vision API — fallback for unconfigured Donut types
    3. Local Donut / Tesseract — opt-in via env flags
    """

    MODEL = "claude-sonnet-4-6"
    MAX_TOKENS = 1024
    DONUT_MODEL = os.getenv("DONUT_MODEL_ID", "naver-clova-ix/donut-base-finetuned-docvqa")
    USE_LOCAL_DONUT = os.getenv("USE_LOCAL_DONUT", "false").strip().lower() in {"1", "true", "yes", "on"}
    USE_LOCAL_TESSERACT = os.getenv("USE_LOCAL_TESSERACT", "false").strip().lower() in {"1", "true", "yes", "on"}

    _donut_processor = None
    _donut_model = None
    _donut_device = None
    _donut_model_id = None

    HF_INFERENCE_URL = "https://api-inference.huggingface.co/models/{model_id}"

    def __init__(self):
        self._client: anthropic.AsyncAnthropic | None = None

    def _get_client(self) -> anthropic.AsyncAnthropic:
        if self._client is None:
            api_key = settings.ANTHROPIC_API_KEY or ""
            if not api_key or api_key == "your-anthropic-api-key-here":
                raise ValueError(
                    "ANTHROPIC_API_KEY is not configured. "
                    "Set it in backend/.env to enable AI document extraction."
                )
            self._client = anthropic.AsyncAnthropic(api_key=api_key)
        return self._client

    async def extract(
        self,
        document_type: str,
        file_bytes: bytes,
        media_type: str,
    ) -> dict:
        """
        Extract structured data from a document image or PDF.

        Returns dict with keys:
          extracted (bool)
          document_type (str)
          data (dict) — on success
          reason (str) — on skip or failure
          message (str) — human-readable description
        """

        normalized_document_type = DOCUMENT_TYPE_ALIASES.get(document_type, document_type)

        if normalized_document_type in SYSTEM_GENERATED_TYPES:
            return {
                "extracted": False,
                "reason": "system_generated",
                "message": "This document type is generated by the system "
                           "and does not require extraction.",
            }

        if normalized_document_type not in EXTRACTION_PROMPTS:
            return {
                "extracted": False,
                "reason": "unknown_type",
                "message": f"No extraction prompt defined for '{document_type}'.",
            }

        # ── Tesseract OCR ──
        _tesseract_methods = {
            "driving_license": self._extract_driving_license_tesseract,
            "rc": self._extract_rc_tesseract,
            "insurance": self._extract_insurance_tesseract,
            "puc": self._extract_puc_tesseract,
            "fitness": self._extract_fitness_tesseract,
            "tax_receipt": self._extract_tax_receipt_tesseract,
            "permit": self._extract_permit_tesseract,
            "aadhaar": self._extract_aadhaar_tesseract,
            "passbook": self._extract_passbook_tesseract,
        }
        tesseract_fn = _tesseract_methods.get(normalized_document_type)
        if tesseract_fn is not None:
            try:
                tess_result = tesseract_fn(
                    file_bytes=file_bytes,
                    media_type=media_type,
                    document_type=document_type,
                )
                return tess_result
            except Exception as e:
                logger.warning("Tesseract extraction failed for %s: %s", document_type, str(e))
                return {
                    "extracted": False,
                    "document_type": document_type,
                    "reason": "tesseract_error",
                    "message": str(e),
                }

        return {
            "extracted": False,
            "document_type": document_type,
            "reason": "unsupported_type",
            "message": f"No extractor available for '{document_type}'.",
        }

    # ─────────────────────────────────────────────────────────
    # HF Inference API — Donut DocVQA (remote, no local model)
    # ─────────────────────────────────────────────────────────

    async def _extract_with_hf_api(
        self, document_type: str, file_bytes: bytes, media_type: str,
    ) -> dict:
        """Call HF Inference API with Donut DocVQA model for each field question."""
        questions = DONUT_FIELD_QUESTIONS.get(document_type)
        if not questions:
            return {
                "extracted": False,
                "reason": "hf_unsupported_type",
                "message": f"Donut prompts are not configured for '{document_type}'.",
            }

        # Convert PDF first page to image if needed
        if media_type == "application/pdf":
            try:
                images = self._pdf_to_images(file_bytes, max_pages=1, dpi=220)
                if not images:
                    return {"extracted": False, "reason": "no_pages", "message": "No pages in PDF."}
                file_bytes = images[0]
            except Exception as e:
                return {"extracted": False, "reason": "pdf_error", "message": str(e)}

        model_id = settings.HF_MODEL_ID
        api_url = self.HF_INFERENCE_URL.format(model_id=model_id)
        hf_key = settings.HF_API_KEY

        image_b64 = base64.b64encode(file_bytes).decode()

        data: dict[str, str | None] = {}

        async with httpx.AsyncClient(timeout=60.0) as client:
            for field, question in questions.items():
                try:
                    payload = {
                        "inputs": {
                            "image": image_b64,
                            "question": question,
                        },
                    }
                    resp = await client.post(
                        api_url,
                        json=payload,
                        headers={"Authorization": f"Bearer {hf_key}"},
                    )

                    if resp.status_code == 503:
                        # Model is loading — wait and retry once
                        body = resp.json()
                        wait_time = min(body.get("estimated_time", 20), 30)
                        logger.info("HF model loading, waiting %.0fs for %s...", wait_time, field)
                        import asyncio
                        await asyncio.sleep(wait_time)
                        resp = await client.post(
                            api_url,
                            json=payload,
                            headers={"Authorization": f"Bearer {hf_key}"},
                        )

                    if resp.status_code != 200:
                        logger.warning("HF API error for field %s: %s %s", field, resp.status_code, resp.text[:200])
                        data[field] = None
                        continue

                    result = resp.json()

                    # HF DocVQA returns: {"answer": "..."} or [{"answer": "...", "score": ...}]
                    answer = None
                    if isinstance(result, dict):
                        answer = result.get("answer")
                    elif isinstance(result, list) and result:
                        answer = result[0].get("answer")

                    answer = self._normalize_donut_field(field, answer)
                    data[field] = answer

                except Exception as e:
                    logger.warning("HF API call failed for field %s: %s", field, str(e))
                    data[field] = None

        if any(v is not None for v in data.values()):
            return {
                "extracted": True,
                "data": data,
            }

        return {
            "extracted": False,
            "reason": "hf_no_fields",
            "message": "HF Donut API could not extract any fields from this document.",
            "data": data,
        }

    def _extract_with_donut(self, document_type: str, file_bytes: bytes, media_type: str) -> dict:
        questions = DONUT_FIELD_QUESTIONS.get(document_type)
        if not questions:
            return {
                "extracted": False,
                "reason": "donut_unsupported_type",
                "message": f"Donut prompts are not configured for '{document_type}'.",
            }

        # PPT/PPTX: extract text from slides then parse using existing rules.
        if media_type in PRESENTATION_MIME_TYPES:
            text = self._extract_text_from_presentation(file_bytes)
            if not text:
                return {
                    "extracted": False,
                    "reason": "presentation_parse_error",
                    "message": "Could not read text from presentation file.",
                }
            data = self._parse_text_by_document_type(document_type, text)
            if any(v is not None for v in data.values()):
                return {
                    "extracted": True,
                    "data": data,
                }
            return {
                "extracted": False,
                "reason": "no_fields",
                "message": "Could not extract expected fields from presentation text.",
                "data": data,
            }

        page_images: list[bytes] = []
        if media_type == "application/pdf":
            try:
                page_images = self._pdf_to_images(file_bytes, max_pages=3, dpi=220)
            except Exception as e:
                return {
                    "extracted": False,
                    "reason": "pdf_conversion_error",
                    "message": f"Could not read PDF: {str(e)}",
                }
        elif media_type in {"image/jpeg", "image/png", "image/webp", "image/heic"}:
            page_images = [file_bytes]
        else:
            return {
                "extracted": False,
                "reason": "donut_unsupported_media",
                "message": f"Unsupported media type for Donut: {media_type}",
            }

        if not page_images:
            return {
                "extracted": False,
                "reason": "no_pages",
                "message": "No pages found to process.",
            }

        try:
            processor, model, device = self._get_donut_runtime()
        except Exception as e:
            return {
                "extracted": False,
                "reason": "donut_unavailable",
                "message": f"Donut model could not be initialized: {str(e)}",
            }

        data: dict[str, str | None] = {k: None for k in questions.keys()}
        for field, question in questions.items():
            answer: str | None = None
            for page in page_images:
                answer = self._ask_donut_question(
                    image_bytes=page,
                    question=question,
                    processor=processor,
                    model=model,
                    device=device,
                )
                answer = self._normalize_donut_field(field, answer)
                if answer:
                    break
            data[field] = answer

        if any(v is not None for v in data.values()):
            return {
                "extracted": True,
                "data": data,
            }

        return {
            "extracted": False,
            "reason": "no_fields",
            "message": "Donut could not detect expected fields.",
            "data": data,
        }

    def _extract_text_from_presentation(self, file_bytes: bytes) -> str:
        try:
            from pptx import Presentation
        except Exception:
            return ""

        try:
            prs = Presentation(io.BytesIO(file_bytes))
        except Exception:
            return ""

        chunks: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text:
                    chunks.append(text)
        return "\n".join(chunks)

    def _parse_text_by_document_type(self, document_type: str, text: str) -> dict:
        if document_type == "driving_license":
            return self._parse_driving_license_text(text)
        if document_type == "rc":
            return self._parse_rc_text(text)
        if document_type == "insurance":
            return self._parse_insurance_text(text)
        if document_type == "puc":
            return self._parse_puc_text(text)
        if document_type == "fitness":
            return self._parse_fitness_text(text)
        if document_type == "tax_receipt":
            return self._parse_tax_receipt_text(text)
        if document_type == "permit":
            return self._parse_permit_text(text)
        return {}

    def _get_donut_runtime(self):
        try:
            import torch
            from transformers import DonutProcessor, VisionEncoderDecoderModel
        except Exception as e:
            raise RuntimeError(
                "Missing Donut dependencies. Install torch and transformers."
            ) from e

        model_id = self.DONUT_MODEL
        if (
            self.__class__._donut_processor is not None
            and self.__class__._donut_model is not None
            and self.__class__._donut_model_id == model_id
        ):
            return (
                self.__class__._donut_processor,
                self.__class__._donut_model,
                self.__class__._donut_device,
            )

        if hasattr(torch.backends, "mps") and torch.backends.mps.is_available():
            device = "mps"
        elif torch.cuda.is_available():
            device = "cuda"
        else:
            device = "cpu"

        processor = DonutProcessor.from_pretrained(model_id)
        model = VisionEncoderDecoderModel.from_pretrained(model_id)
        model.to(device)
        model.eval()

        self.__class__._donut_processor = processor
        self.__class__._donut_model = model
        self.__class__._donut_device = device
        self.__class__._donut_model_id = model_id
        return processor, model, device

    def _ask_donut_question(self, image_bytes: bytes, question: str, processor, model, device) -> str | None:
        try:
            from PIL import Image
            import torch
        except Exception:
            return None

        try:
            image = Image.open(io.BytesIO(image_bytes)).convert("RGB")
            task_prompt = f"<s_docvqa><s_question>{question}</s_question><s_answer>"
            decoder_input_ids = processor.tokenizer(
                task_prompt,
                add_special_tokens=False,
                return_tensors="pt",
            ).input_ids.to(device)

            pixel_values = processor(image, return_tensors="pt").pixel_values.to(device)

            with torch.no_grad():
                outputs = model.generate(
                    pixel_values,
                    decoder_input_ids=decoder_input_ids,
                    max_length=256,
                    pad_token_id=processor.tokenizer.pad_token_id,
                    eos_token_id=processor.tokenizer.eos_token_id,
                    bad_words_ids=[[processor.tokenizer.unk_token_id]],
                    use_cache=True,
                )

            decoded = processor.batch_decode(outputs, skip_special_tokens=False)[0]
            if "<s_answer>" in decoded:
                decoded = decoded.split("<s_answer>", 1)[1]
            if "</s_answer>" in decoded:
                decoded = decoded.split("</s_answer>", 1)[0]
            cleaned = re.sub(r"<[^>]+>", " ", decoded)
            cleaned = re.sub(r"\s+", " ", cleaned).strip()
            return cleaned or None
        except Exception:
            return None

    def _normalize_donut_field(self, field: str, value: str | None) -> str | None:
        if value is None:
            return None

        cleaned = re.sub(r"\s+", " ", value).strip(" :-")
        if not cleaned:
            return None

        if cleaned.lower() in {
            "n/a", "na", "none", "null", "not found", "unknown", "not available", "-", "--", "—"
        }:
            return None

        if field in {
            "license_number", "registration_number", "vehicle_number", "permit_number",
            "certificate_number", "receipt_number", "policy_number", "pan_number", "aadhaar_number",
        }:
            compact = re.sub(r"\s+", "", cleaned).upper()
            return compact or None

        if "date" in field or field in {"valid_until", "validity_date"}:
            normalized = self._normalize_date(cleaned)
            return normalized or cleaned

        if field == "fuel_type":
            fv = cleaned.lower()
            if "diesel" in fv:
                return "Diesel"
            if "petrol" in fv:
                return "Petrol"
            if "cng" in fv:
                return "CNG"
            if "electric" in fv:
                return "Electric"
            if "lpg" in fv:
                return "LPG"

        return cleaned

    # ─────────────────────────────────────────────────────────
    # Shared OCR helpers (used by all Tesseract extractors)
    # ─────────────────────────────────────────────────────────

    def _preprocess_for_ocr(self, image: "Image.Image") -> "Image.Image":  # type: ignore[name-defined]
        """
        Apply a standard preprocessing pipeline to improve Tesseract accuracy.

        Steps:
        1. Convert to grayscale (already done by callers but safe to repeat).
        2. Upscale so the long edge is at least 2400 px — Tesseract accuracy
           degrades sharply below ~150 dpi equivalent.
        3. Otsu / adaptive binarization via a gentle threshold so the background
           is white and ink is black.
        4. Unsharp-mask sharpening to enhance character edges.
        5. Autocontrast to maximise the tonal range.
        """
        from PIL import Image, ImageOps, ImageFilter
        import io

        image = image.convert("L")

        # ── 1. Upscale ──
        w, h = image.size
        min_long_edge = 2400
        if max(w, h) < min_long_edge:
            scale = min_long_edge / max(w, h)
            image = image.resize((int(w * scale), int(h * scale)), Image.LANCZOS)

        # ── 2. Autocontrast before binarisation ──
        image = ImageOps.autocontrast(image, cutoff=2)

        # ── 3. Sharpen ──
        image = image.filter(ImageFilter.UnsharpMask(radius=1.5, percent=180, threshold=3))

        # ── 4. Binarise with Otsu's threshold using a histogram approach ──
        try:
            import numpy as np
            arr = np.array(image)
            # Compute histogram
            hist, bins = np.histogram(arr.flatten(), bins=256, range=(0, 256))
            # Otsu's method: find threshold that maximises inter-class variance
            total = arr.size
            sum_total = float(np.dot(np.arange(256), hist))
            sum_bg, weight_bg = 0.0, 0
            best_var, threshold = 0.0, 128
            for t in range(256):
                weight_bg += hist[t]
                if weight_bg == 0:
                    continue
                weight_fg = total - weight_bg
                if weight_fg == 0:
                    break
                sum_bg += t * hist[t]
                mean_bg = sum_bg / weight_bg
                mean_fg = (sum_total - sum_bg) / weight_fg
                var = weight_bg * weight_fg * (mean_bg - mean_fg) ** 2
                if var > best_var:
                    best_var, threshold = var, t
            image = image.point(lambda p: 255 if p >= threshold else 0)
        except ImportError:
            # numpy not available — fallback to fixed threshold
            image = image.point(lambda p: 255 if p >= 128 else 0)

        return image

    def _run_tesseract_multipass(self, image: "Image.Image") -> str:  # type: ignore[name-defined]
        """
        Run Tesseract in multiple PSM modes and languages, returning combined text.
        Uses PSM 6 (uniform block), 4 (single column), and 11 (sparse text).
        Tries English first, then Hindi+English for Indian documents.
        Returns the longest non-empty result from the first successful pass.
        """
        import pytesseract

        # Locate tesseract on Windows if needed
        if os.name == "nt":
            current = getattr(pytesseract.pytesseract, "tesseract_cmd", "tesseract")
            if not os.path.exists(current):
                for win_path in [
                    r"C:\Program Files\Tesseract-OCR\tesseract.exe",
                    r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
                ]:
                    if os.path.exists(win_path):
                        pytesseract.pytesseract.tesseract_cmd = win_path
                        break

        results: list[str] = []
        for psm in ("6", "4", "11", "3"):
            for lang in ("eng", "eng+hin"):
                try:
                    t = pytesseract.image_to_string(
                        image, lang=lang,
                        config=f"--oem 3 --psm {psm}",
                    ) or ""
                    if t.strip():
                        results.append(t)
                except Exception:
                    pass
            if results:
                break  # stop at first PSM mode that yields text

        # Also try to get any remaining text from other modes and append
        for psm in ("4", "11"):
            for lang in ("eng",):
                try:
                    t = pytesseract.image_to_string(
                        image, lang=lang,
                        config=f"--oem 3 --psm {psm}",
                    ) or ""
                    if t.strip() and t not in results:
                        results.append(t)
                except Exception:
                    pass

        return "\n".join(results)

    def _extract_driving_license_tesseract(
        self,
        file_bytes: bytes,
        media_type: str,
        document_type: str,
    ) -> dict:
        """Extract key DL fields via Tesseract OCR (offline path)."""
        try:
            from PIL import Image
            import io
            import pytesseract  # noqa: F401 — triggers the import check
        except Exception as e:
            return {
                "extracted": False,
                "reason": "tesseract_unavailable",
                "message": f"Tesseract OCR dependencies are missing: {str(e)}",
            }

        image_pages: list[bytes] = []
        if media_type == "application/pdf":
            try:
                image_pages = self._pdf_to_images(file_bytes, max_pages=3, dpi=300)
            except Exception as e:
                return {
                    "extracted": False,
                    "reason": "pdf_conversion_error",
                    "message": f"Could not read PDF: {str(e)}",
                }
        else:
            image_pages = [file_bytes]

        text_chunks: list[str] = []
        try:
            for page_bytes in image_pages:
                raw_image = Image.open(io.BytesIO(page_bytes)).convert("L")
                image = self._preprocess_for_ocr(raw_image)
                text_chunks.append(self._run_tesseract_multipass(image))
            text = "\n".join(chunk for chunk in text_chunks if chunk)
        except Exception as e:
            return {
                "extracted": False,
                "reason": "ocr_error",
                "message": f"Tesseract OCR failed: {str(e)}",
            }

        data = self._parse_driving_license_text(text)
        if not any(v is not None for v in data.values()):
            return {
                "extracted": False,
                "reason": "no_fields",
                "message": "Could not detect license number or dates from the document.",
                "data": data,
            }

        return {
            "extracted": True,
            "document_type": document_type,
            "data": data,
        }

    def _extract_rc_tesseract(
        self,
        file_bytes: bytes,
        media_type: str,
        document_type: str,
    ) -> dict:
        """Extract key RC fields via Tesseract OCR (offline path)."""
        try:
            from PIL import Image
            import io
            import pytesseract  # noqa: F401
        except Exception as e:
            return {"extracted": False, "reason": "tesseract_unavailable", "message": str(e)}

        if media_type == "application/pdf":
            try:
                file_bytes, _ = self._pdf_to_image(file_bytes)
            except Exception as e:
                return {"extracted": False, "reason": "pdf_conversion_error", "message": str(e)}

        try:
            image = self._preprocess_for_ocr(Image.open(io.BytesIO(file_bytes)))
            text = self._run_tesseract_multipass(image)
        except Exception as e:
            return {"extracted": False, "reason": "ocr_error", "message": str(e)}

        data = self._parse_rc_text(text)
        if not any(v is not None for v in data.values()):
            return {"extracted": False, "reason": "no_fields", "message": "Could not detect RC fields from the document.", "data": data}

        return {"extracted": True, "document_type": document_type, "data": data}

    def _extract_insurance_tesseract(
        self,
        file_bytes: bytes,
        media_type: str,
        document_type: str,
    ) -> dict:
        """Extract key insurance fields via Tesseract OCR (offline path)."""
        try:
            from PIL import Image
            import io
            import pytesseract  # noqa: F401
        except Exception as e:
            return {"extracted": False, "reason": "tesseract_unavailable", "message": str(e)}

        if media_type == "application/pdf":
            try:
                file_bytes, _ = self._pdf_to_image(file_bytes)
            except Exception as e:
                return {"extracted": False, "reason": "pdf_conversion_error", "message": str(e)}

        try:
            image = self._preprocess_for_ocr(Image.open(io.BytesIO(file_bytes)))
            text = self._run_tesseract_multipass(image)
        except Exception as e:
            return {"extracted": False, "reason": "ocr_error", "message": str(e)}

        data = self._parse_insurance_text(text)
        if not any(v is not None for v in data.values()):
            return {"extracted": False, "reason": "no_fields", "message": "Could not detect insurance fields from the document.", "data": data}

        return {"extracted": True, "document_type": document_type, "data": data}

    def _extract_puc_tesseract(
        self,
        file_bytes: bytes,
        media_type: str,
        document_type: str,
    ) -> dict:
        """Extract key PUC fields via Tesseract OCR (offline path)."""
        try:
            from PIL import Image
            import io
            import pytesseract  # noqa: F401
        except Exception as e:
            return {"extracted": False, "reason": "tesseract_unavailable", "message": str(e)}

        if media_type == "application/pdf":
            try:
                file_bytes, _ = self._pdf_to_image(file_bytes)
            except Exception as e:
                return {"extracted": False, "reason": "pdf_conversion_error", "message": str(e)}

        try:
            image = self._preprocess_for_ocr(Image.open(io.BytesIO(file_bytes)))
            text = self._run_tesseract_multipass(image)
        except Exception as e:
            return {"extracted": False, "reason": "ocr_error", "message": str(e)}

        data = self._parse_puc_text(text)
        if not any(v is not None for v in data.values()):
            return {"extracted": False, "reason": "no_fields", "message": "Could not detect PUC fields from the document.", "data": data}

        return {"extracted": True, "document_type": document_type, "data": data}

    def _extract_fitness_tesseract(
        self,
        file_bytes: bytes,
        media_type: str,
        document_type: str,
    ) -> dict:
        """Extract key fitness certificate fields via Tesseract OCR."""
        try:
            from PIL import Image
            import io
            import pytesseract  # noqa: F401
        except Exception as e:
            return {"extracted": False, "reason": "tesseract_unavailable", "message": str(e)}

        if media_type == "application/pdf":
            try:
                file_bytes, _ = self._pdf_to_image(file_bytes)
            except Exception as e:
                return {"extracted": False, "reason": "pdf_conversion_error", "message": str(e)}

        try:
            image = self._preprocess_for_ocr(Image.open(io.BytesIO(file_bytes)))
            text = self._run_tesseract_multipass(image)
        except Exception as e:
            return {"extracted": False, "reason": "ocr_error", "message": str(e)}

        data = self._parse_fitness_text(text)
        if not any(v is not None for v in data.values()):
            return {"extracted": False, "reason": "no_fields", "message": "Could not detect fitness fields from the document.", "data": data}

        return {"extracted": True, "document_type": document_type, "data": data}

    def _extract_tax_receipt_tesseract(
        self,
        file_bytes: bytes,
        media_type: str,
        document_type: str,
    ) -> dict:
        """Extract key road-tax receipt fields via Tesseract OCR."""
        try:
            from PIL import Image
            import io
            import pytesseract  # noqa: F401
        except Exception as e:
            return {"extracted": False, "reason": "tesseract_unavailable", "message": str(e)}

        if media_type == "application/pdf":
            try:
                file_bytes, _ = self._pdf_to_image(file_bytes)
            except Exception as e:
                return {"extracted": False, "reason": "pdf_conversion_error", "message": str(e)}

        try:
            image = self._preprocess_for_ocr(Image.open(io.BytesIO(file_bytes)))
            text = self._run_tesseract_multipass(image)
        except Exception as e:
            return {"extracted": False, "reason": "ocr_error", "message": str(e)}

        data = self._parse_tax_receipt_text(text)
        if not any(v is not None for v in data.values()):
            return {"extracted": False, "reason": "no_fields", "message": "Could not detect tax receipt fields from the document.", "data": data}

        return {"extracted": True, "document_type": document_type, "data": data}

    def _extract_permit_tesseract(
        self,
        file_bytes: bytes,
        media_type: str,
        document_type: str,
    ) -> dict:
        """Extract key permit fields via Tesseract OCR."""
        try:
            from PIL import Image
            import io
            import pytesseract  # noqa: F401
        except Exception as e:
            return {"extracted": False, "reason": "tesseract_unavailable", "message": str(e)}

        if media_type == "application/pdf":
            try:
                file_bytes, _ = self._pdf_to_image(file_bytes)
            except Exception as e:
                return {"extracted": False, "reason": "pdf_conversion_error", "message": str(e)}

        try:
            image = self._preprocess_for_ocr(Image.open(io.BytesIO(file_bytes)))
            text = self._run_tesseract_multipass(image)
        except Exception as e:
            return {"extracted": False, "reason": "ocr_error", "message": str(e)}

        data = self._parse_permit_text(text)
        if not any(v is not None for v in data.values()):
            return {"extracted": False, "reason": "no_fields", "message": "Could not detect permit fields from the document.", "data": data}

        return {"extracted": True, "document_type": document_type, "data": data}

    # ── Aadhaar OCR extractor ─────────────────────────────────────────────────

    def _extract_aadhaar_tesseract(
        self,
        file_bytes: bytes,
        media_type: str,
        document_type: str,
    ) -> dict:
        """Extract Aadhaar fields (DOB, gender, address, number) via OCR. Name is intentionally excluded."""
        try:
            from PIL import Image
            import io
            import pytesseract  # noqa: F401
        except Exception as e:
            return {"extracted": False, "reason": "tesseract_unavailable", "message": str(e)}

        if media_type == "application/pdf":
            try:
                file_bytes, _ = self._pdf_to_image(file_bytes)
            except Exception as e:
                return {"extracted": False, "reason": "pdf_conversion_error", "message": str(e)}

        try:
            image = self._preprocess_for_ocr(Image.open(io.BytesIO(file_bytes)))
            text = self._run_tesseract_multipass(image)
        except Exception as e:
            return {"extracted": False, "reason": "ocr_error", "message": str(e)}

        data = self._parse_aadhaar_text(text)
        if not any(v is not None for v in data.values()):
            return {
                "extracted": False,
                "reason": "no_fields",
                "message": "Could not detect Aadhaar fields from the document.",
                "data": data,
            }
        return {"extracted": True, "document_type": document_type, "data": data}

    def _parse_aadhaar_text(self, text: str) -> dict:
        """Parse OCR text from an Aadhaar card image."""
        lines = [ln.strip() for ln in (text or "").splitlines() if ln.strip()]
        joined = " ".join(lines)

        # ── Aadhaar number: 4-4-4 digit groups ──
        aadhaar_number = None
        m = re.search(r"\b(\d{4})\s+(\d{4})\s+(\d{4})\b", joined)
        if m:
            aadhaar_number = m.group(1) + m.group(2) + m.group(3)

        # ── Date of birth ──
        dob = None
        dob_patterns = [
            r"(?:DOB|Date\s*of\s*Birth|D\.O\.B)[:\s]+(\d{2}[\/\-]\d{2}[\/\-]\d{4})",
            r"(?:DOB|Date\s*of\s*Birth|D\.O\.B)[:\s]+(\d{4})",
            r"\b(\d{2}[\/\-]\d{2}[\/\-]\d{4})\b",
        ]
        for pat in dob_patterns:
            m = re.search(pat, joined, re.IGNORECASE)
            if m:
                raw = m.group(1).replace("-", "/")
                # Normalise to DD/MM/YYYY
                if re.match(r"\d{4}$", raw):
                    dob = raw  # year only
                elif re.match(r"\d{2}/\d{2}/\d{4}", raw):
                    dob = raw
                else:
                    dob = raw
                break

        # ── Gender ──
        gender = None
        if re.search(r"\bFEMALE\b", joined, re.IGNORECASE):
            gender = "Female"
        elif re.search(r"\bMALE\b", joined, re.IGNORECASE):
            gender = "Male"
        elif re.search(r"\bTRANSGENDER\b", joined, re.IGNORECASE):
            gender = "Transgender"

        return {
            "date_of_birth": dob,
            "gender": gender,
            "aadhaar_number": aadhaar_number,
        }

    def _extract_next_line_after_label(self, lines: list, label_patterns: list, max_chars: int = 120) -> str | None:
        """Extract the value from the line AFTER a matching label line (TN RC smart card layout)."""
        for i, line in enumerate(lines):
            for pattern in label_patterns:
                if re.search(pattern, line, re.IGNORECASE):
                    # Check if value is on same line after the label
                    same_line = re.sub(pattern, "", line, flags=re.IGNORECASE).strip(" :-")
                    if same_line and len(same_line) > 1:
                        return same_line[:max_chars].strip()
                    # Otherwise take next non-empty line
                    for j in range(i + 1, min(i + 3, len(lines))):
                        candidate = lines[j].strip()
                        if candidate:
                            return candidate[:max_chars]
        return None

    def _parse_rc_text(self, text: str) -> dict:
        cleaned = re.sub(r"[\t\r]+", " ", text or "")
        cleaned_upper = cleaned.upper()
        lines = [ln.strip() for ln in cleaned.split("\n") if ln.strip()]

        registration_number = self._extract_registration_number(cleaned_upper, lines)

        # TN RC smart card: labels are on one line, values on the NEXT line
        owner_name = self._extract_next_line_after_label(
            lines,
            [r"OWNER\s*NAME", r"NAME\s*OF\s*OWNER"],
            max_chars=120,
        )
        owner_name = self._clean_owner_name(owner_name)
        if owner_name and not self._is_plausible_owner_name(owner_name):
            owner_name = None
        if not owner_name:
            owner_name = self._extract_owner_name_from_lines(lines)
        vehicle_class = self._extract_rc_labeled_text(
            cleaned,
            [
                r"VEHICLE\s*CLASS",
                r"CLASS\s*OF\s*VEHICLE",
            ],
            max_chars=60,
        )

        # Fuel is usually on same line as label on TN RC (left side)
        fuel_type = self._normalize_fuel_type(
            self._extract_rc_labeled_text(
                cleaned,
                [r"FUEL\s*USED", r"FUEL\s*TYPE", r"^FUEL"],
                max_chars=40,
            ) or self._extract_next_line_after_label(
                lines,
                [r"^FUEL$", r"FUEL\s*USED", r"FUEL\s*TYPE"],
                max_chars=40,
            )
        )

        # Engine/Motor Number and Chassis Number: next-line extraction
        engine_number = self._extract_next_line_after_label(
            lines,
            [r"ENGINE\s*/\s*MOTOR\s*NUMBER", r"ENGINE\s*MOTOR\s*NUMBER", r"ENGINE\s*NUMBER", r"ENGINE\s*NO"],
            max_chars=60,
        )
        # Validate: must be alphanumeric code, not a name
        if engine_number and not re.search(r"[A-Z0-9]{4,}", engine_number.upper()):
            engine_number = None
        if engine_number and re.fullmatch(r"[A-Z ]+", engine_number.upper()):
            engine_number = None

        chassis_number = self._extract_next_line_after_label(
            lines,
            [r"CHASSIS\s*NUMBER", r"CHASSIS\s*NO"],
            max_chars=60,
        )
        # Validate: must be long alphanumeric VIN, not a name
        if chassis_number and re.fullmatch(r"[A-Z ]+", chassis_number.upper()):
            chassis_number = None
        if chassis_number and len(chassis_number) < 8:
            chassis_number = None

        issue_date = self._extract_date_from_line_context(
            lines,
            [r"DATE\s*OF\s*REGN", r"DATE\s*OF\s*REGISTRATION", r"REGN\.?\s*DATE"],
            pick="first",
        )
        validity_date = self._extract_date_from_line_context(
            lines,
            [r"REGN\.?\s*VALIDITY", r"VALID\s*UPTO", r"VALIDITY", r"VALID\s*TILL", r"RC\s*VALID\s*TILL"],
            pick="latest",
        )

        if not issue_date or not validity_date:
            all_dates = self._sort_dates_unique(self._extract_all_dates(cleaned))
            if all_dates:
                if not issue_date:
                    issue_date = all_dates[0]
                if not validity_date:
                    validity_date = all_dates[-1]

        return {
            "registration_number": registration_number,
            "owner_name": owner_name,
            "vehicle_class": vehicle_class,
            "fuel_type": fuel_type,
            "engine_number": engine_number,
            "chassis_number": chassis_number,
            "issue_date": issue_date,
            "validity_date": validity_date,
        }

    def _parse_insurance_text(self, text: str) -> dict:
        cleaned = re.sub(r"[\t\r]+", " ", text or "")
        cleaned_upper = cleaned.upper()
        lines = [ln.strip() for ln in cleaned.split("\n") if ln.strip()]

        policy_number = self._extract_rc_labeled_text(
            cleaned,
            [
                r"POLICY\s*NO\.?",
                r"CERTIFICATE\s*NO\.?",
                r"POLICY\s*NUMBER",
            ],
            max_chars=60,
        )

        insurer_name = self._extract_rc_labeled_text(
            cleaned,
            [
                r"NAME\s*OF\s*INSURER",
                r"INSURANCE\s*COMPANY",
                r"INSURED\s*BY",
            ],
            max_chars=120,
        )

        vehicle_number = self._extract_registration_number(cleaned_upper, lines)

        issue_date = self._extract_date_from_line_context(
            lines,
            [
                r"POLICY\s*ISSUE\s*DATE",
                r"DATE\s*OF\s*ISSUE",
                r"COMMENCEMENT\s*DATE",
            ],
            pick="first",
        )

        expiry_date = self._extract_date_from_line_context(
            lines,
            [
                r"POLICY\s*EXPIRY\s*DATE",
                r"DATE\s*OF\s*EXPIRY",
                r"VALID\s*TILL",
                r"EXPIRY\s*DATE",
            ],
            pick="latest",
        )

        if not issue_date or not expiry_date:
            all_dates = self._sort_dates_unique(self._extract_all_dates(cleaned))
            if all_dates:
                if not issue_date:
                    issue_date = all_dates[0]
                if not expiry_date:
                    expiry_date = all_dates[-1]

        return {
            "policy_number": policy_number,
            "insurer_name": insurer_name,
            "issue_date": issue_date,
            "expiry_date": expiry_date,
            "vehicle_number": vehicle_number,
        }

    def _parse_puc_text(self, text: str) -> dict:
        cleaned = re.sub(r"[\t\r]+", " ", text or "")
        cleaned_upper = cleaned.upper()
        lines = [ln.strip() for ln in cleaned.split("\n") if ln.strip()]

        certificate_number = self._extract_rc_labeled_text(
            cleaned,
            [
                r"PUC\s*CERT\.?\s*NO\.?",
                r"CERTIFICATE\s*NO\.?",
                r"TEST\s*CERTIFICATE\s*NO\.?",
                r"CERT\s*NO\.?",
                r"PUC\s*NO\.?",
            ],
            max_chars=60,
        )

        vehicle_number = self._extract_registration_number(cleaned_upper, lines)

        test_date = self._extract_date_from_line_context(
            lines,
            [
                r"DATE\s*OF\s*TEST",
                r"TEST\s*DATE",
                r"TESTED\s*ON",
            ],
            pick="first",
        )

        valid_until = self._extract_date_from_line_context(
            lines,
            [
                r"VALID\s*UPTO",
                r"VALID\s*TILL",
                r"CERTIFICATE\s*VALID\s*TILL",
            ],
            pick="latest",
        )

        if not test_date or not valid_until:
            all_dates = self._sort_dates_unique(self._extract_all_dates(cleaned))
            if all_dates:
                if not test_date:
                    test_date = all_dates[0]
                if not valid_until:
                    valid_until = all_dates[-1]

        reading_values = self._extract_puc_readings(cleaned)

        return {
            "certificate_number": certificate_number,
            "vehicle_number": vehicle_number,
            "test_date": test_date,
            "valid_until": valid_until,
            "reading_values": reading_values,
        }

    def _parse_fitness_text(self, text: str) -> dict:
        """
        Parse OCR text from a Tamil Nadu E-Fee Receipt / Fitness Certificate.

        Handles both the traditional FC layout and the ONLINEPERMITFEE portal
        E-Fee Receipt format (which carries Fitness Validity, Tax Paid Upto,
        Chassis No, and Grand Total alongside the standard vehicle/owner fields).
        """
        cleaned = re.sub(r"[\t\r]+", " ", text or "")

        # ── Raw OCR text corrections (applied before any parsing) ────────────
        # 1. Receipt Date: OCR drops leading digit of two-digit day.
        #    "Receipt Date: 8-Feb-2025" or "08-Feb-2025" → try to fix using the
        #    transaction number encoded date (handled later after extraction).
        # 2. Chassis No prefix IBI/MBI/MB| → MB1 (Ashok Leyland WMI = MB1).
        #    OCR confuses M→I, 1→I giving "IBI"; or M→M, 1→I giving "MBI".
        cleaned = re.sub(
            r'\b(IBI|MBI|MB\|)([A-Z0-9]{14})\b',
            r'MB1\2',
            cleaned.upper(),
        )

        cleaned_upper = cleaned.upper()
        lines = [ln.strip() for ln in cleaned.split("\n") if ln.strip()]

        # ── Vehicle number ───────────────────────────────────────────────────
        vehicle_number = self._extract_registration_number(cleaned_upper, lines)

        # ── Owner name ───────────────────────────────────────────────────────
        owner_name = self._extract_rc_labeled_text(
            cleaned,
            [
                r"OWNER\s*NAME",
                r"NAME\s*OF\s*OWNER",
                r"REGISTERED\s*OWNER",
                r"PERMIT\s*HOLDER\s*NAME",
                r"NAME\s*OF\s*THE\s*PERMIT\s*HOLDER",
            ],
            max_chars=120,
        )
        owner_name = self._clean_owner_name(owner_name)
        if owner_name and not self._is_plausible_owner_name(owner_name):
            owner_name = None
        if not owner_name:
            owner_name = self._extract_owner_name_from_lines(lines)

        # ── Fitness validity (the future date FC is valid until) ─────────────
        # Label: "Fitness Validity", "Valid Upto", "FC Valid Till", etc.
        fitness_validity = self._extract_date_from_line_context(
            lines,
            [
                r"FITNESS\s*VALIDITY",
                r"FITNESS\s*VALID\s*(?:UPTO|TILL|TO)",
                r"FC\s*VALID\s*(?:TILL|UPTO|TO)",
                r"VALID\s*UPTO",
                r"VALID\s*TILL",
                r"EXPIRY\s*DATE",
                r"VALIDITY",
            ],
            pick="latest",
        )

        # ── Tax Paid Upto ────────────────────────────────────────────────────
        tax_paid_upto = self._extract_date_from_line_context(
            lines,
            [
                r"TAX\s*PAID\s*(?:UPTO|TILL|TO|UP\s*TO)",
                r"ROAD\s*TAX\s*(?:PAID|VALID)\s*(?:UPTO|TILL|TO|UP\s*TO)",
            ],
            pick="latest",
        )

        # ── Receipt / issue date ──────────────────────────────────────────────
        receipt_date = self._extract_date_from_line_context(
            lines,
            [
                r"RECEIPT\s*DATE",
                r"DATE\s*OF\s*(?:ISSUE|RECEIPT|PAYMENT)",
                r"ISSUE\s*DATE",
            ],
            pick="first",
        )

        # ── Chassis number: 17-char alphanumeric VIN ─────────────────────────
        chassis_no = self._extract_rc_labeled_text(
            cleaned,
            [r"CHASSIS\s*NO\.?", r"CHASSIS\s*NUMBER"],
            max_chars=25,
        )
        # Validate / clean: must be ≥10 chars, alphanumeric
        if chassis_no:
            chassis_no = re.sub(r"\s+", "", chassis_no.upper())
            if len(chassis_no) < 10 or not re.match(r"^[A-Z0-9]+$", chassis_no):
                chassis_no = None
        if not chassis_no:
            # Scan lines for a standalone 17-char VIN-like token
            for ln in lines:
                m = re.search(r"\b([A-Z0-9]{17})\b", ln.upper())
                if m:
                    candidate = m.group(1)
                    # Must have both letters and digits
                    if re.search(r"[A-Z]", candidate) and re.search(r"\d", candidate):
                        chassis_no = candidate
                        break
        # ── Chassis OCR correction: IBI → MB1 (Ashok Leyland prefix)
        # OCR commonly reads M→I, B→B, 1→I yielding "IBI" instead of "MB1"
        if chassis_no and chassis_no.startswith("IBI"):
            chassis_no = "MB1" + chassis_no[3:]

        # ── Transaction / Receipt number ──────────────────────────────────────
        transaction_no = self._extract_rc_labeled_text(
            cleaned,
            [
                r"TRANSACTION\s*NO\.?(?:\s*/\s*RECEIPT\s*NO\.?)?",
                r"RECEIPT\s*NO\.?",
                r"TRANSACTION\s*(?:ID|NUMBER)",
            ],
            max_chars=80,
        )
        if transaction_no:
            transaction_no = re.sub(r"\s+", "", transaction_no.upper())
        if not transaction_no:
            # Pattern: TN + digits + letters + digits, optionally paired with /
            m = re.search(
                r"(TN\d{6}[A-Z]\d+(?:/TN\d{6}[A-Z]\d+)?)",
                cleaned_upper,
            )
            if m:
                transaction_no = m.group(1)

        # ── Extract encoded date from transaction number ───────────────────────
        # TN E-Fee format: TN + YYMMDD + letter + serial
        # e.g. TN250228P9553804 → 2025-02-28
        # Also scan raw OCR for partial matches when TN prefix is garbled
        # (OCR sometimes drops T, producing just N250228... or "7250228...")
        tn_encoded_date: str | None = None
        _tn_scan_targets = [transaction_no or "", cleaned_upper]
        for _src in _tn_scan_targets:
            _m = re.search(r"(?:TN|N)(\d{2})(\d{2})(\d{2})[A-Z]", _src)
            if _m:
                yy, mm, dd = _m.groups()
                try:
                    _dt = datetime.strptime(f"{dd}/{mm}/20{yy}", "%d/%m/%Y")
                    # Sanity: year must be between 2020 and 2035
                    if 2020 <= _dt.year <= 2035:
                        tn_encoded_date = _dt.strftime("%d/%m/%Y")
                        break
                except ValueError:
                    pass

        # ── Cross-validate receipt_date against encoded transaction date ───────
        if tn_encoded_date:
            if receipt_date:
                try:
                    rd_dt = datetime.strptime(receipt_date, "%d/%m/%Y")
                    tn_dt = datetime.strptime(tn_encoded_date, "%d/%m/%Y")
                    # If year+month match but day differs → trust transaction number
                    if rd_dt.year == tn_dt.year and rd_dt.month == tn_dt.month and rd_dt.day != tn_dt.day:
                        receipt_date = tn_encoded_date
                except ValueError:
                    pass
            else:
                receipt_date = tn_encoded_date

        # ── Office name ───────────────────────────────────────────────────────
        office = self._extract_rc_labeled_text(
            cleaned,
            [r"OFFICE\s*[:\-]?", r"ISSUING\s*(?:OFFICE|AUTHORITY|RTO)"],
            max_chars=80,
        )
        if not office:
            # Look for lines containing "RTO" or "REGIONAL TRANSPORT"
            for ln in lines:
                if re.search(r"\bRTO\b|\bREGIONAL\s*TRANSPORT\b", ln, re.IGNORECASE):
                    candidate = ln.strip()
                    # Remove leading "Office :" prefix if present
                    candidate = re.sub(r"^.*?OFFICE\s*[:\-]\s*", "", candidate, flags=re.IGNORECASE).strip()
                    if candidate and len(candidate) > 3:
                        office = candidate.upper()
                        break

        # ── Grand total ───────────────────────────────────────────────────────
        grand_total = None
        m = re.search(
            r"GRAND\s*TOTAL\s*(?:\([^)]*\))?\s*[:\-]?\s*([\d,]+(?:\.\d+)?)",
            cleaned, re.IGNORECASE,
        )
        if m:
            grand_total = m.group(1).replace(",", "")
        if not grand_total:
            # Fallback: look for "1100/-" style pattern on a GRAND TOTAL line
            for ln in lines:
                if re.search(r"GRAND\s*TOTAL", ln, re.IGNORECASE):
                    amt = re.search(r"([\d,]+)(?:/\s*-|\s*Rs|\s*INR)", ln)
                    if amt:
                        grand_total = amt.group(1).replace(",", "")
                        break

        # ── Traditional FC certificate number ────────────────────────────────
        certificate_number = self._extract_rc_labeled_text(
            cleaned,
            [
                r"CERTIFICATE\s*NO\.?",
                r"FITNESS\s*CERT\.?\s*NO\.?",
                r"FC\s*NO\.?",
            ],
            max_chars=60,
        )

        # ── Fallback: populate missing dates from all dates found ─────────────
        if not fitness_validity or not receipt_date:
            all_dates = self._sort_dates_unique(self._extract_all_dates(cleaned))
            if all_dates:
                if not receipt_date:
                    receipt_date = all_dates[0]
                if not fitness_validity:
                    fitness_validity = all_dates[-1]

        return {
            "certificate_number": certificate_number,
            "vehicle_number": vehicle_number,
            "owner_name": owner_name,
            "fitness_validity": fitness_validity,
            "tax_paid_upto": tax_paid_upto,
            "receipt_date": receipt_date,
            "chassis_no": chassis_no,
            "transaction_no": transaction_no,
            "office": office,
            "grand_total": grand_total,
            # Keep legacy aliases for backward compat
            "issue_date": receipt_date,
            "expiry_date": fitness_validity,
        }

    def _parse_tax_receipt_text(self, text: str) -> dict:
        cleaned = re.sub(r"[\t\r]+", " ", text or "")
        cleaned_upper = cleaned.upper()
        lines = [ln.strip() for ln in cleaned.split("\n") if ln.strip()]

        receipt_number = self._extract_rc_labeled_text(
            cleaned,
            [
                r"RECEIPT\s*NO\.?",
                r"TRANSACTION\s*NO\.?",
                r"CHALLAN\s*NO\.?",
                r"TAX\s*RECEIPT\s*NO\.?",
            ],
            max_chars=60,
        )
        vehicle_number = self._extract_registration_number(cleaned_upper, lines)
        tax_amount = self._extract_tax_amount(cleaned)
        valid_until = self._extract_date_from_line_context(
            lines,
            [
                r"VALID\s*UPTO",
                r"TAX\s*VALID\s*TILL",
                r"VALIDITY",
            ],
            pick="latest",
        )

        if not valid_until:
            all_dates = self._sort_dates_unique(self._extract_all_dates(cleaned))
            if all_dates:
                valid_until = all_dates[-1]

        return {
            "receipt_number": receipt_number,
            "vehicle_number": vehicle_number,
            "tax_amount": tax_amount,
            "valid_until": valid_until,
        }

    def _parse_permit_text(self, text: str) -> dict:
        cleaned = re.sub(r"[\t\r]+", " ", text or "")
        cleaned_upper = cleaned.upper()
        lines = [ln.strip() for ln in cleaned.split("\n") if ln.strip()]

        permit_number = self._extract_permit_number(cleaned_upper)
        date_of_approval = self._extract_date_from_line_context(
            lines,
            [r"DATE\s*OF\s*APPROVAL", r"APPROVAL\s*DATE"],
            pick="first",
        )

        permit_type = None
        if re.search(r"NATIONAL\s*PERMIT", cleaned_upper):
            permit_type = "NATIONAL PERMIT"
        elif re.search(r"FRESH\s*PERMIT", cleaned_upper):
            permit_type = "FRESH PERMIT"

        permit_holder_name = self._extract_permit_holder_name(lines)
        father_name = self._extract_father_name(lines)
        address = self._extract_address_block(lines)

        registration_number = self._extract_registration_number(cleaned_upper, lines)
        type_of_vehicle = self._extract_rc_labeled_text(
            cleaned,
            [r"TYPE\s*OF\s*VEHICLE", r"VEHICLE\s*TYPE", r"CLASS\s*OF\s*VEHICLE"],
            max_chars=120,
        )
        # Reject garbled type_of_vehicle: must contain a known vehicle class keyword
        if type_of_vehicle:
            _veh_keywords = {"GOODS", "CARRIER", "PASSENGER", "TANKER", "BUS", "TAXI",
                             "TRUCK", "TRAILER", "LMV", "HGV", "HMV", "LCV", "MGV"}
            _words = set(re.sub(r"[^A-Z ]", "", type_of_vehicle.upper()).split())
            if not _veh_keywords & _words:
                type_of_vehicle = None
        # Fallback: permit uses two-column numbered layout — label and value on separate lines.
        # Directly search for known Indian vehicle type strings anywhere in the text.
        if not type_of_vehicle:
            _known_veh_types = [
                ("Goods Carrier",         r"\bGoods\s+Carri(?:er|age)\b"),
                ("Passenger Carrier",     r"\bPassenger\s+Carri(?:er|age)\b"),
                ("Stage Carriage",        r"\bStage\s+Carriage\b"),
                ("Contract Carriage",     r"\bContract\s+Carriage\b"),
                ("Private Service Vehicle", r"\bPrivate\s+Service\s+Vehicle\b"),
                ("Motor Cab",             r"\bMotor\s+Cab\b"),
                ("Maxi Cab",              r"\bMaxi\s+Cab\b"),
                ("Tanker",                r"\bTanker\b"),
                ("LMV",                   r"\bLMV\b"),
            ]
            for _canonical, _pat in _known_veh_types:
                if re.search(_pat, cleaned, re.IGNORECASE):
                    type_of_vehicle = _canonical
                    break

        valid_from = self._extract_date_from_line_context(
            lines,
            [r"FROM\s*[:\-]", r"VALID\s*FROM"],
            pick="first",
        )
        valid_to = self._extract_date_from_line_context(
            lines,
            [r"TO\s*[:\-]", r"VALID\s*TO", r"VALID\s*UPTO", r"EXPIRY\s*DATE"],
            pick="latest",
        )

        # Prefer explicit validity pair extraction from OCR text (From:- ... To:-).
        # This is more reliable than independent label/date matching when both
        # dates appear on the same noisy line.
        explicit_valid_from, explicit_valid_to = self._extract_validity_range(cleaned)
        if explicit_valid_from:
            valid_from = explicit_valid_from
        if explicit_valid_to:
            valid_to = explicit_valid_to

        # Also try cross-line validity range: OCR sometimes splits From/To across two lines.
        if not explicit_valid_to or valid_to == valid_from:
            _date_tok = re.compile(
                r"(\d{1,2}[\s\-/][A-Za-z]{3,9}[\s\-/]\d{4})"
            )
            for _i, _ln in enumerate(lines):
                if not re.search(r"\bFrom\b", _ln, re.IGNORECASE):
                    continue
                _from_dates = _date_tok.findall(_ln)
                if not _from_dates:
                    continue
                # Check this line and the next two for a To: date
                for _j in range(_i, min(_i + 3, len(lines))):
                    if not re.search(r"\bTo\b", lines[_j], re.IGNORECASE):
                        continue
                    _to_dates = _date_tok.findall(lines[_j])
                    if not _to_dates:
                        continue
                    _s = self._normalize_date(_from_dates[0])
                    _e = self._normalize_date(_to_dates[-1])
                    if self._is_reasonable_validity_range(_s, _e):
                        valid_from = _s
                        valid_to = _e
                    break
                break

        states_uts_covered = self._extract_rc_labeled_text(
            cleaned,
            [r"VALID\s*FOR", r"STATES?\s*/?\s*UTS?\s*COVERED", r"AREA\s*OF\s*OPERATION"],
            max_chars=160,
        )

        issued_date = self._extract_date_from_line_context(
            lines,
            [r"DATE\s*[:\-]", r"DATE\s*OF\s*ISSUE", r"ISSUED\s*ON"],
            pick="latest",
        )

        # Authorization number: e.g. TN2023-NP/AUTH-7395A
        _auth_match = re.search(
            r"\b([A-Z]{2}\d{4}[-/ ]*NP[-/ ]*AUTH[-/ ]*\d{3,5}[A-Z]?)\b",
            cleaned_upper,
        )
        authorization_no = re.sub(r"\s+", "", _auth_match.group(1)).upper() if _auth_match else None

        # Authorization validity: "From: 22-Feb-2023  To: 21-Feb-2024" on auth-validity line
        authorization_validity_from: str | None = None
        authorization_validity_to: str | None = None
        for _ln in lines:
            if re.search(r"AUTHORI[ZS]ATION\s*VALIDITY", _ln, re.IGNORECASE):
                _vf, _vt = self._extract_validity_range(_ln)
                if _vf or _vt:
                    authorization_validity_from = _vf
                    authorization_validity_to = _vt
                break

        issuing_authority = self._extract_issuing_authority(lines)
        # Reject clearly garbled issuing_authority values
        if issuing_authority:
            _ia_words = issuing_authority.split()
            _short_noise = sum(1 for w in _ia_words if len(w) <= 2)
            if (
                not _ia_words
                or (_ia_words[0] and _ia_words[0][0].isdigit())
                or re.search(r'["\'\*\|\\<>]', issuing_authority)
                or (_ia_words and _short_noise > len(_ia_words) // 2)
            ):
                issuing_authority = None

        # Backward-compatible fields used by current UI/form mapping.
        issue_date = date_of_approval or valid_from or issued_date
        expiry_date = valid_to
        route_area = states_uts_covered
        vehicle_number = registration_number

        all_dates = self._sort_dates_unique(self._extract_all_dates(cleaned))
        if not valid_from or not valid_to:
            if all_dates:
                if not valid_from:
                    valid_from = all_dates[0]
                if not valid_to:
                    valid_to = all_dates[-1]
                issue_date = issue_date or valid_from
                expiry_date = expiry_date or valid_to

        # Last resort: if valid_to is still same as valid_from pick the latest date
        # in the document that falls AFTER valid_from (any duration is better than same date).
        if valid_from and (not valid_to or valid_to == valid_from):
            try:
                _vf_dt = datetime.strptime(valid_from, "%d/%m/%Y")
                for _d in reversed(all_dates):
                    try:
                        _dd = datetime.strptime(_d, "%d/%m/%Y")
                        if _dd > _vf_dt:
                            valid_to = _d
                            expiry_date = _d
                            break
                    except ValueError:
                        continue
            except ValueError:
                pass

        return {
            "permit_number": permit_number,
            "date_of_approval": date_of_approval,
            "permit_type": permit_type,
            "permit_holder_name": permit_holder_name,
            "father_name": father_name,
            "address": address,
            "registration_number": registration_number,
            "type_of_vehicle": type_of_vehicle,
            "valid_from": valid_from,
            "valid_to": valid_to,
            "states_uts_covered": states_uts_covered,
            "issued_date": issued_date,
            "issuing_authority": issuing_authority,
            "authorization_no": authorization_no,
            "authorization_validity_from": authorization_validity_from,
            "authorization_validity_to": authorization_validity_to,
            "vehicle_number": vehicle_number,
            "route_area": route_area,
            "issue_date": issue_date,
            "expiry_date": expiry_date,
        }

    def _extract_permit_number(self, cleaned_upper: str) -> str | None:
        patterns = [
            re.compile(r"\b([A-Z]{2}\d{4}-NP-\d{4}[A-Z]?)\b"),
            re.compile(r"\b([A-Z]{2}\d{4}\s*[-/]?\s*NP\s*[-/]?\s*\d{4}[A-Z]?)\b"),
        ]
        for pattern in patterns:
            match = pattern.search(cleaned_upper)
            if match:
                return re.sub(r"\s+", "", match.group(1).replace("/", "-")).upper()

        labeled = self._extract_rc_labeled_text(
            cleaned_upper,
            [r"PERMIT\s*NO\.?", r"AUTHORI[ZS]ATION\s*NO\.?", r"PERMIT\s*NUMBER"],
            max_chars=60,
        )
        if not labeled:
            return None
        return re.sub(r"\s+", "", labeled.replace("/", "-")).upper()

    def _extract_validity_range(self, text: str) -> tuple[str | None, str | None]:
        """Return the From/To date pair with the LONGEST span found in text.

        A permit document may contain multiple From/To pairs (e.g. permit validity
        AND authorization validity).  We always want the longest span — the 5-year
        national permit — not the shorter 1-year authorization validity.
        """
        date_token = r"(\d{1,2}(?:[\s\-/])[A-Za-z]{3,9}(?:[\s\-/])\d{4})"

        patterns = [
            re.compile(rf"From\s*:-\s*{date_token}\s+To\s*:-\s*{date_token}", re.IGNORECASE),
            re.compile(rf"From\s*[:-=]+\s*{date_token}\s+To\s*[:-=]+\s*{date_token}", re.IGNORECASE),
            re.compile(rf"From\s*[:\-=]+\s*{date_token}[\s\S]{{0,120}}?To\s*[:\-=]+\s*{date_token}", re.IGNORECASE),
            re.compile(rf"From[^0-9]*{date_token}[^0-9]*To[^0-9]*{date_token}", re.IGNORECASE),
        ]

        best_start: str | None = None
        best_end: str | None = None
        best_span = 0

        for pattern in patterns:
            for match in pattern.finditer(text):
                start = self._normalize_date(match.group(1))
                end = self._normalize_date(match.group(2))
                if not self._is_reasonable_validity_range(start, end):
                    continue
                try:
                    span = (datetime.strptime(end, "%d/%m/%Y") - datetime.strptime(start, "%d/%m/%Y")).days  # type: ignore[arg-type]
                except ValueError:
                    continue
                if span > best_span:
                    best_span = span
                    best_start = start
                    best_end = end

        if best_start:
            return best_start, best_end

        # Line-level fallback for OCR where label punctuation is corrupted.
        date_pair = re.compile(r"(\d{1,2}(?:[\s\-/])[A-Za-z]{3,9}(?:[\s\-/])\d{4})")
        for line in text.splitlines():
            if not re.search(r"\bFROM\b", line, re.IGNORECASE) or not re.search(r"\bTO\b", line, re.IGNORECASE):
                continue
            tokens = date_pair.findall(line)
            if len(tokens) < 2:
                continue
            start = self._normalize_date(tokens[0])
            end = self._normalize_date(tokens[-1])
            if not self._is_reasonable_validity_range(start, end):
                continue
            try:
                span = (datetime.strptime(end, "%d/%m/%Y") - datetime.strptime(start, "%d/%m/%Y")).days  # type: ignore[arg-type]
            except ValueError:
                continue
            if span > best_span:
                best_span = span
                best_start = start
                best_end = end

        return best_start, best_end

    def _is_reasonable_validity_range(self, start: str | None, end: str | None) -> bool:
        if not start or not end:
            return False

        try:
            start_dt = datetime.strptime(start, "%d/%m/%Y")
            end_dt = datetime.strptime(end, "%d/%m/%Y")
        except ValueError:
            return False

        if start_dt.year < 2000 or end_dt.year < 2000:
            return False

        return end_dt > start_dt

    def _extract_text_after_label_from_lines(self, lines: list[str], label_patterns: list[str]) -> str | None:
        compiled = [re.compile(p, re.IGNORECASE) for p in label_patterns]
        for i, line in enumerate(lines):
            if not any(p.search(line) for p in compiled):
                continue

            inline = re.sub(r"^.*?:", "", line).strip(" :.-")
            if inline and len(inline) > 2 and not any(p.search(inline) for p in compiled):
                return inline

            for nxt in lines[i + 1 : i + 4]:
                probe = nxt.strip(" :.-")
                if not probe:
                    continue
                if any(p.search(probe) for p in compiled):
                    continue
                if re.search(r"^\d{1,2}[\-/][A-Za-z]{3}[\-/]\d{2,4}$", probe):
                    continue
                return probe
        return None

    def _extract_permit_holder_name(self, lines: list[str]) -> str | None:
        label_pattern = re.compile(r"NAME\s*OF\s*THE\s*PERMIT\s*HOLDER|PERMIT\s*HOLDER\s*NAME", re.IGNORECASE)

        for i, line in enumerate(lines):
            if not label_pattern.search(line):
                continue

            inline = re.sub(r"^\s*\d+\.?\s*", "", line).strip()
            inline = re.sub(r".*?(NAME\s*OF\s*THE\s*PERMIT\s*HOLDER|PERMIT\s*HOLDER\s*NAME)\s*[:\-]?", "", inline, flags=re.IGNORECASE).strip(" :.-")
            inline = self._strip_following_numbered_label(inline)
            if self._is_plausible_permit_holder_name(inline):
                return inline

            for nxt in lines[i + 1 : i + 4]:
                probe = re.sub(r"^\s*\d+\.?\s*", "", nxt).strip(" :.-")
                if not probe:
                    continue
                if re.search(r"\bFATHER'?S\s*NAME\b|\bADDRESS\b|\bVALID\b|\bTYPE\s*OF\s*VEHICLE\b|\bREG(?:ISTRATION)?\s*NO\b", probe, re.IGNORECASE):
                    break
                probe = self._strip_following_numbered_label(probe)
                if self._is_plausible_permit_holder_name(probe):
                    return probe

        return None

    def _is_plausible_permit_holder_name(self, value: str | None) -> bool:
        if not value:
            return False
        cleaned = re.sub(r"\s+", " ", value).strip(" :.-")
        if not cleaned:
            return False
        if cleaned.upper() in {"NA", "N/A", "NIL", "NULL"}:
            return False
        if re.search(r"\bFATHER'?S\s*NAME\b|\bADDRESS\b|\bVALID\b|\bPERMIT\s*NO\b", cleaned, re.IGNORECASE):
            return False
        if re.search(r"\d{3,}", cleaned):
            return False
        return len(cleaned) >= 3

    def _strip_following_numbered_label(self, value: str | None) -> str | None:
        if not value:
            return None
        cleaned = re.sub(r"\s+", " ", value).strip(" :.-")
        # OCR can merge sections like:
        # "2. Name Of The Permit Holder 3. Father's Name NA"
        # Strip leading numbering and cut off when the next numbered label starts.
        cleaned = re.sub(r"^\d+\.\s*", "", cleaned)
        cleaned = re.split(r"(?:^|\s)\d+\.\s*[A-Za-z]", cleaned, maxsplit=1)[0].strip(" :.-")
        return cleaned or None

    def _extract_father_name(self, lines: list[str]) -> str | None:
        label_patterns = [
            re.compile(r"FATHER'?S\s*NAME", re.IGNORECASE),
            re.compile(r"\bS/O\b", re.IGNORECASE),
            re.compile(r"\bD/O\b", re.IGNORECASE),
            re.compile(r"\bW/O\b", re.IGNORECASE),
        ]

        for i, line in enumerate(lines):
            if not any(p.search(line) for p in label_patterns):
                continue

            inline = re.sub(r"^\s*\d+\.?\s*", "", line).strip()
            inline = re.sub(r".*?(FATHER'?S\s*NAME|\bS/O\b|\bD/O\b|\bW/O\b)\s*[:\-]?", "", inline, flags=re.IGNORECASE).strip(" :.-")
            inline = self._strip_following_numbered_label(inline)
            if inline and inline.upper() in {"NA", "N/A", "NIL", "NULL", "NOT AVAILABLE"}:
                return None
            if self._is_plausible_father_name(inline):
                return inline

            for nxt in lines[i + 1 : i + 4]:
                probe = re.sub(r"^\s*\d+\.?\s*", "", nxt).strip(" :.-")
                if not probe:
                    continue
                if re.search(r"\bADDRESS\b|\bTYPE\s*OF\s*VEHICLE\b|\bREG(?:ISTRATION)?\s*NO\b|\bVALID\b|\bFROM\b|\bTO\b", probe, re.IGNORECASE):
                    break
                if any(p.search(probe) for p in label_patterns):
                    continue
                probe = self._strip_following_numbered_label(probe)
                if probe and probe.upper() in {"NA", "N/A", "NIL", "NULL", "NOT AVAILABLE"}:
                    return None
                if self._is_plausible_father_name(probe):
                    return probe

        return None

    def _is_plausible_father_name(self, value: str | None) -> bool:
        if not value:
            return False
        cleaned = re.sub(r"\s+", " ", value).strip(" :.-")
        if not cleaned:
            return False
        if cleaned.upper() in {"NA", "N/A", "NIL", "NULL", "NOT AVAILABLE"}:
            return False
        if re.search(r"\d", cleaned):
            return False
        if re.search(r"\bADDRESS\b|\bHOTEL\b|\bSHOP\b|\bROAD\b|\bSTREET\b|\bNAGAR\b|\bBYPASS\b|\bTAMIL\s*NADU\b", cleaned, re.IGNORECASE):
            return False
        words = [w for w in cleaned.split(" ") if w]
        if len(words) < 1 or len(words) > 6:
            return False
        return True

    def _extract_address_block(self, lines: list[str]) -> str | None:
        for i, line in enumerate(lines):
            if not re.search(r"\bADDRESS\b", line, re.IGNORECASE):
                continue

            collected: list[str] = []

            # OCR often emits: "4. Address HOTEL ..." on one line.
            inline = re.sub(r"^\s*\d+\.?\s*", "", line).strip()
            inline = re.sub(r".*?\bADDRESS\b\s*[:\-]?", "", inline, flags=re.IGNORECASE).strip(" :.-")
            inline = self._strip_following_numbered_label(inline)
            if inline and not re.search(r"\b(?:TYPE\s*OF\s*VEHICLE|FROM\s*[:\-]|TO\s*[:\-]|VALID\s*FOR|DATE\s*OF\s*APPROVAL|PERMIT\s*NO|REGN\s*NO|REGISTRATION\s*NO|ISSUING\s*AUTHORITY|REGIONAL\s*TRANSPORT\s*AUTHORITY|RTA|RTO|AUTHORI[ZS]ATION\s*NO|NATIONAL\s*PERMIT|FRESH\s*PERMIT|MOTOR\s*VEHICLE|DEPARTMENT)\b", inline, re.IGNORECASE):
                collected.append(inline)

            for nxt in lines[i + 1 : i + 7]:
                probe = nxt.strip(" :.-")
                if not probe:
                    if collected:
                        break
                    continue
                if re.search(r"\b(?:TYPE\s*OF\s*VEHICLE|FROM\s*[:\-]|TO\s*[:\-]|VALID\s*FOR|DATE\s*OF\s*APPROVAL|PERMIT\s*NO|REGN\s*NO|REGISTRATION\s*NO|ISSUING\s*AUTHORITY|REGIONAL\s*TRANSPORT\s*AUTHORITY|RTA|RTO|AUTHORI[ZS]ATION\s*NO|NATIONAL\s*PERMIT|FRESH\s*PERMIT|MOTOR\s*VEHICLE|DEPARTMENT)\b", probe, re.IGNORECASE):
                    break
                collected.append(probe)

            if collected:
                return ", ".join(collected)

        return None

    def _extract_issuing_authority(self, lines: list[str]) -> str | None:
        labeled = self._extract_text_after_label_from_lines(
            lines,
            [r"ISSUING\s*AUTHORITY", r"AUTHORITY", r"REGIONAL\s*TRANSPORT\s*AUTHORITY", r"TRANSPORT\s*AUTHORITY", r"RTA", r"RTO"],
        )
        if labeled:
            return labeled

        for line in lines:
            if re.search(r"REGIONAL\s+TRANSPORT\s+AUTHORITY", line, re.IGNORECASE):
                return re.sub(r"\s+", " ", line).strip(" :.-")
            if re.search(r"\bRTA\b|\bRTO\b", line, re.IGNORECASE) and re.search(r"AUTHORITY|OFFICE|CHENNAI|TRANSPORT", line, re.IGNORECASE):
                return re.sub(r"\s+", " ", line).strip(" :.-")
        return None

    def _extract_puc_readings(self, text: str) -> dict | None:
        fields = [
            ("CO% Vol", r"CO\s*%\s*VOL"),
            ("HC (PPM)", r"HC\s*\(?\s*PPM\s*\)?"),
            ("CO2%", r"CO2\s*%"),
            ("Idle CO%", r"IDLE\s*CO\s*%"),
            ("High Idle HC", r"HIGH\s*IDLE\s*HC"),
        ]
        out: dict[str, str] = {}

        for label, rx in fields:
            pattern = re.compile(rx + r"\s*[:=\-]?\s*([0-9]+(?:\.[0-9]+)?)", re.IGNORECASE)
            match = pattern.search(text)
            if match:
                out[label] = match.group(1)

        return out or None

    def _extract_tax_amount(self, text: str) -> str | None:
        labels = [
            r"TAX\s*AMOUNT",
            r"AMOUNT\s*PAID",
            r"TAX\s*PAID",
            r"TOTAL\s*AMOUNT",
        ]
        for label in labels:
            pattern = re.compile(label + r"\s*[:=\-]?\s*([₹Rs\.\s]*\d[\d,]*(?:\.\d{1,2})?)", re.IGNORECASE)
            match = pattern.search(text)
            if not match:
                continue
            raw = re.sub(r"\s+", " ", match.group(1)).strip()
            if "₹" in raw:
                return raw
            if re.search(r"\bRS\.?\b", raw, re.IGNORECASE):
                return raw
            return raw
        return None

    def _extract_registration_number(self, cleaned_upper: str, lines: list[str]) -> str | None:
        # Permit OCR often prints "Regn No/Manuf. year ... TN72CB4704 / 2023".
        permit_line_pattern = re.compile(
            r"(?:REG(?:ISTRATION)?\s*(?:NO|NUMBER)?|REGN\s*NO\.?|VEHICLE\s*NO\.?|MANUF\.?\s*YEAR[^\n]{0,60})[^A-Z0-9]*([A-Z]{2}\s*\d{1,2}\s*[A-Z]{1,3}\s*\d{3,4})",
            re.IGNORECASE,
        )
        permit_line_match = permit_line_pattern.search(cleaned_upper)
        if permit_line_match:
            normalized = self._normalize_registration_number(permit_line_match.group(1))
            if normalized:
                return normalized

        labeled_pattern = re.compile(
            r"(?:REG(?:ISTRATION)?\s*(?:NO|NUMBER)?|REG\.\s*NO\.?|VEHICLE\s*NO\.?|REGN\s*NO\.?)\s*[:\-]?\s*([A-Z]{2}\s*\d{1,2}\s*[A-Z]{1,3}\s*\d{3,4})",
            re.IGNORECASE,
        )
        labeled_match = labeled_pattern.search(cleaned_upper)
        if labeled_match:
            normalized = self._normalize_registration_number(labeled_match.group(1))
            if normalized:
                return normalized

        top_pattern = re.compile(r"\b[A-Z]{2}\s*\d{1,2}\s*[A-Z]{1,3}\s*\d{3,4}\b")
        for ln in lines:
            match = top_pattern.search(ln.upper())
            if match:
                normalized = self._normalize_registration_number(match.group(0))
                if normalized:
                    return normalized
        return None

    def _normalize_registration_number(self, raw: str | None) -> str | None:
        if not raw:
            return None

        value = re.sub(r"[^A-Z0-9]", "", raw.upper())
        match = re.fullmatch(r"([A-Z]{2})(\d{1,2})([A-Z]{1,3})(\d{3,4})", value)
        if not match:
            return None

        state, district, series, number = match.groups()

        # OCR on Tamil Nadu RCs commonly misreads TN as IN (low-contrast) or FN (T→F).
        if state in {"IN", "FN"}:
            state = "TN"

        return f"{state}{district}{series}{number}"

    def _is_plausible_owner_name(self, value: str | None) -> bool:
        if not value:
            return False
        candidate = self._clean_owner_name(value)
        if not candidate:
            return False
        if len(candidate) < 4:
            return False
        if len(candidate) == 1:
            return False
        if re.search(r"\d", candidate):
            return False
        if candidate.upper() in {"OWNER", "SERIAL", "NAME", "OWNER NAME"}:
            return False
        words = [w for w in candidate.split(" ") if w]
        return len(words) >= 2

    def _extract_owner_name_from_lines(self, lines: list[str]) -> str | None:
        # On many Indian RC cards, owner name appears on the line immediately
        # after the "Owner Name" label and before address/father-name lines.
        for i, line in enumerate(lines):
            if not re.search(r"OWNER\s*NAME|NAME\s*OF\s*OWNER", line, re.IGNORECASE):
                continue

            # Try inline value first, then next few lines.
            inline = re.sub(r".*?(OWNER\s*NAME|NAME\s*OF\s*OWNER)\s*[:\-]?", "", line, flags=re.IGNORECASE).strip(" :.-")
            inline = self._clean_owner_name(inline)
            if self._is_plausible_owner_name(inline):
                return inline

            for nxt in lines[i + 1 : i + 4]:
                probe = self._clean_owner_name(nxt)
                if re.search(r"S/O|D/O|W/O|SON|DAUGHTER|WIFE|ADDRESS|REGN|CHASSIS|ENGINE", probe, re.IGNORECASE):
                    continue
                if self._is_plausible_owner_name(probe):
                    return probe

        return None

    def _clean_owner_name(self, value: str | None) -> str | None:
        if not value:
            return None

        candidate = re.sub(r"\s+", " ", value).strip(" :.-()")
        if not candidate:
            return None

        # Remove common OCR garbage tokens that appear around owner lines.
        words = [w for w in candidate.split(" ") if w]
        while words and words[0].lower() in {"if", "f", "i", "of"}:
            words.pop(0)
        # Strip trailing single chars that are not letters (OCR artifacts like '(')
        while words and (len(words[-1]) == 1 and not words[-1].isalpha()):
            words.pop()
        while words and len(words[-1]) == 1 and words[-1].islower():
            words.pop()

        cleaned = " ".join(words).strip(" :.-()")
        return cleaned or None

    def _extract_rc_labeled_text(self, text: str, labels: list[str], max_chars: int = 80) -> str | None:
        for label in labels:
            pattern = re.compile(label + r"\s*[:\-]?\s*([^\n]{1," + str(max_chars) + r"})", re.IGNORECASE)
            match = pattern.search(text)
            if not match:
                continue

            value = match.group(1).strip(" :.-")
            # Trim value if OCR captured the next inline label.
            value = re.split(
                r"\s{2,}|\b(?:REG(?:ISTRATION)?\s*(?:NO|NUMBER)?|FUEL\s*(?:TYPE|USED)|ENGINE\s*(?:NO|NUMBER)|CHASSIS\s*(?:NO|NUMBER)|DATE|VALID(?:ITY|\s*TILL|\s*UPTO))\b",
                value,
                maxsplit=1,
                flags=re.IGNORECASE,
            )[0].strip(" :.-")
            return value or None
        return None

    def _normalize_fuel_type(self, value: str | None) -> str | None:
        if not value:
            return None
        text = value.strip().lower()
        if "diesel" in text:
            return "Diesel"
        if "petrol" in text:
            return "Petrol"
        if "cng" in text:
            return "CNG"
        if "electric" in text or "ev" == text:
            return "Electric"
        if "lpg" in text:
            return "LPG"
        return value.strip()

    def _parse_driving_license_text(self, text: str) -> dict:
        cleaned = re.sub(r"[\t\r]+", " ", text or "")
        lines = [ln.strip() for ln in cleaned.split("\n") if ln.strip()]

        # ── Words to exclude from name matching ──
        _HEADER_WORDS = {
            # Document headers
            "indian", "union", "driving", "licence", "license", "issued",
            "government", "tamil", "nadu", "kerala", "karnataka", "andhra",
            "pradesh", "maharashtra", "telangana", "state", "transport",
            "authority", "department", "republic", "india",
            # Field labels
            "date", "birth", "validity", "expiry", "issue", "issued",
            "blood", "group", "grup", "address", "son", "daughter", "wife",
            "husband", "father", "mother", "name", "number", "class",
            "vehicle", "ref", "type", "badge", "pin", "code", "signature",
            "photo", "holder", "valid", "from", "upto", "till", "the", "of",
        }

        # ══════════════════════════════════════════════
        # 1. LICENSE NUMBER
        # ══════════════════════════════════════════════
        label_pattern = re.compile(
            r"(?:LICENSE\s*(?:NO|NUMBER|#)?|LICENCE\s*(?:NO|NUMBER|#)?|DL\s*(?:NO|NUMBER|#)?|DRIVING\s+LICEN[CS]E\s*(?:NO|NUMBER)?)\s*[:\-]?\s*([A-Z0-9\-/\s]{6,32})",
            re.IGNORECASE,
        )
        label_match = label_pattern.search(cleaned)
        labeled_license = self._normalize_license_token(label_match.group(1)) if label_match else None
        if labeled_license and not self._is_probable_driving_license_number(labeled_license):
            labeled_license = None

        top_candidate = self._extract_driving_license_number_from_lines(lines)

        if not top_candidate:
            token_pattern = re.compile(r"\b[A-Z0-9\-/]{6,25}\b")
            for ln in lines[:8]:
                for tok in token_pattern.findall(ln.upper()):
                    normalized = self._normalize_license_token(tok)
                    if normalized and self._is_probable_driving_license_number(normalized):
                        top_candidate = normalized
                        break
                if top_candidate:
                    break

        license_number = labeled_license or top_candidate

        # Fallback: collapse all whitespace/newlines and retry the DL pattern.
        # Handles cases where OCR splits "TN72" and "20240005499" onto separate lines.
        if not license_number:
            _dl_pat_collapsed = re.compile(
                r"\b([A-Z]{2}\s*[-/]?\s*\d{1,2}\s*[-/]?\s*\d{4}\s*[-/]?\s*\d{4,10})\b",
                re.IGNORECASE,
            )
            collapsed = re.sub(r"\s+", " ", cleaned.upper())
            for match in _dl_pat_collapsed.findall(collapsed):
                norm = self._normalize_license_token(match)
                if norm and self._is_probable_driving_license_number(norm):
                    license_number = norm
                    break

        # Fix truncated state codes: OCR often drops leading letter(s)
        # e.g. "N7220240005499" → "TN7220240005499"
        if not license_number:
            trunc = re.search(r"\b([A-Z]?\d{2}\s*\d{9,13})\b", cleaned.upper())
            if trunc:
                raw = re.sub(r"\s+", "", trunc.group(1))
                # If starts with single letter + 2 digits, probably missing first letter
                if re.match(r"^[A-Z]\d{2}\d{9,13}$", raw):
                    # Try common state codes
                    for prefix in ["T", "K", "A", "M"]:
                        candidate = prefix + raw
                        if self._is_probable_driving_license_number(candidate):
                            license_number = candidate
                            break
                elif re.match(r"^\d{2}\d{9,13}$", raw):
                    # Both letters missing — less reliable, skip
                    pass

        # ══════════════════════════════════════════════
        # 2. DATES — issue, expiry, DOB
        # ══════════════════════════════════════════════
        issue_date = self._extract_date_from_line_context(
            lines,
            [r"(?:I|L)SSUE\s*DATE", r"DATE\s*OF\s*ISSUE", r"ISSUED\s*ON", r"\bDOI\b"],
            pick="first",
        )

        validity_dates = self._extract_labeled_dates(
            cleaned,
            r"(?:VALIDITY\s*(?:\(\s*(?:NT|TR)\s*\))?|VALID\s*(?:UPTO|UNTIL|TILL|TO)|EXPIRY\s*DATE|DATE\s*OF\s*EXPIRY|EXPIRES\s*ON|VALID\s*\(\s*(?:NT|TR)\s*\))",
        )
        for p in [r"VALIDITY", r"VALID\s*(?:UPTO|UNTIL|TILL|TO)", r"EXPIRY", r"\bNT\b", r"\bTR\b"]:
            d = self._extract_date_from_line_context(lines, [p], pick="latest")
            if d:
                validity_dates.append(d)

        dob_dates = self._extract_labeled_dates(cleaned, r"(?:DATE\s*OF\s*BIRTH|DOB|BIRTH)")
        validity_dates = [d for d in validity_dates if d not in dob_dates]
        expiry_date = self._pick_latest_date(validity_dates)

        dob = self._extract_labeled_date(cleaned, r"(?:DATE\s*OF\s*BIRTH|DOB|D\.?O\.?B\.?|BIRTH\s*DATE)")

        # ── DOB age-sanity correction ─────────────────────────────────────────
        # OCR often confuses digits (5↔8, 6↔0, 1↔7).  If the extracted DOB
        # would make the holder younger than 17 at the issue date, try common
        # single-digit substitutions on the year until we find a valid age.
        if dob and issue_date:
            try:
                _dob_day, _dob_mon, _dob_yr_str = dob.split("/")
                _iss_day, _iss_mon, _iss_yr_str = issue_date.split("/")
                _dob_yr = int(_dob_yr_str)
                _iss_yr = int(_iss_yr_str)
                _age_at_issue = _iss_yr - _dob_yr
                if _age_at_issue < 17:
                    # Common OCR digit-swap pairs (wrong → [candidates])
                    _ocr_swaps = {
                        "8": ["5", "6", "3"],
                        "6": ["5", "0"],
                        "0": ["8", "6"],
                        "7": ["1", "4"],
                        "1": ["7"],
                        "3": ["8"],
                    }
                    _yr = _dob_yr_str
                    _corrected = False
                    for _pos in [3, 2]:          # try last digit first, then second-to-last
                        _digit = _yr[_pos]
                        for _sub in _ocr_swaps.get(_digit, []):
                            _new_yr_str = _yr[:_pos] + _sub + _yr[_pos + 1:]
                            _new_age = _iss_yr - int(_new_yr_str)
                            if 17 <= _new_age <= 80:
                                dob = f"{_dob_day}/{_dob_mon}/{_new_yr_str}"
                                _corrected = True
                                break
                        if _corrected:
                            break
            except Exception:
                pass

        if not issue_date:
            issue_from_validity = self._extract_labeled_date(cleaned, r"(?:VALID\s*FROM|VALIDITY\s*FROM)")
            if issue_from_validity:
                issue_date = issue_from_validity

        # Smart date assignment from all found dates
        if not issue_date or not expiry_date:
            all_dates = self._extract_all_dates(cleaned)
            all_dates = [d for d in all_dates if d not in dob_dates]
            if all_dates:
                unique_dates = self._sort_dates_unique(all_dates)
                if not issue_date:
                    issue_date = unique_dates[0]
                if not expiry_date:
                    expiry_date = unique_dates[-1]

        # Sanity check: if issue_date == dob, it's probably wrong
        if issue_date and dob and issue_date == dob:
            all_dates = self._extract_all_dates(cleaned)
            # Remove DOB and find next earliest
            other_dates = [d for d in all_dates if d != dob]
            if other_dates:
                sorted_other = self._sort_dates_unique(other_dates)
                issue_date = sorted_other[0]

        # ══════════════════════════════════════════════
        # 3. HOLDER NAME — label-first, then scoring
        # ══════════════════════════════════════════════
        holder_name = None

        def _clean_candidate(raw: str) -> str:
            """Strip leading/trailing label words from a candidate."""
            words = raw.split()
            while words and words[0].lower() in _HEADER_WORDS:
                words = words[1:]
            while words and words[-1].lower() in _HEADER_WORDS:
                words = words[:-1]
            return " ".join(words)

        def _is_valid_name(text: str) -> bool:
            """Return True if text looks like a real person name."""
            words = text.split()
            # Need at least 1 word of a real length; single-word names accepted only if >= 4 chars
            if not words:
                return False
            if len(text.replace(" ", "")) < 4:
                return False
            if re.search(r"\d", text):
                return False
            # Must contain at least one word longer than 2 chars (rules out "En", "OC" etc.)
            if not any(len(w) > 2 for w in words):
                return False
            avg_len = sum(len(w) for w in words) / len(words)
            if avg_len <= 2:
                return False
            lower = text.lower()
            if any(h in lower for h in [
                "driving", "licence", "license", "government", "issued",
                "indian union", "tamil nadu", "kerala", "karnataka",
                "endorsement", "entitlement", "engine", "chassis",
                "transport", "department", "authority",
            ]):
                return False
            return True

        # ── Strategy 1: NEXT-LINE search (preferred for Indian DLs) ──────────
        # On Indian smart-card DLs "Name" is a column header on one line,
        # and the actual name sits on the ROW BELOW it.
        # We try this FIRST because same-line captures are often garbled
        # OCR fragments from adjacent columns (e.g. "Endorsements").
        label_line_pat = re.compile(
            r"(?:holder'?s?\s*name|name\s*of\s*(?:dl\s*)?holder|dl\s*holder|\bname\b)",
            re.IGNORECASE,
        )
        _so_line_pat = re.compile(
            r"^\s*(?:s/?o|d/?o|w/?o|son|daughter|wife|husband)\b",
            re.IGNORECASE,
        )
        for i, line in enumerate(lines):
            if label_line_pat.search(line):
                # Try lines directly after the label line (up to 3 lines ahead)
                for j in range(i + 1, min(i + 4, len(lines))):
                    # Skip S/O, D/O, W/O lines — those are the parent's name
                    if _so_line_pat.match(lines[j]):
                        continue
                    candidate = _clean_candidate(lines[j])
                    if not candidate or len(candidate.replace(" ", "")) < 4:
                        continue
                    if _is_valid_name(candidate):
                        # Extra guard: don't accept lines that look like field labels
                        if re.search(r"[:/]|\b(no|of|at|the)\b", candidate, re.IGNORECASE):
                            continue
                        holder_name = candidate.title()
                        break
                if holder_name:
                    break

        # ── Strategy 1b: inline "Name: AJAI KUMAR" on the SAME line ─────────
        # Fallback for DLs where OCR places the name on the same line as the label.
        if not holder_name:
            name_label_pat = re.compile(
                r"(?:holder'?s?\s*name|name\s*of\s*(?:dl\s*)?holder|dl\s*holder|name)\s*[:\-]?\s*"
                r"([A-Z][A-Z\s\.]{3,60})",  # min 4 total chars
                re.IGNORECASE,
            )
            for m in name_label_pat.finditer(cleaned):
                raw = m.group(1).strip()
                # Trim at a line break or common delimiters
                raw = re.split(r"\n|Date|DOB|Son|Daughter|Wife|Address|Blood|S/O|D/O|W/O", raw, flags=re.IGNORECASE)[0].strip()
                candidate = _clean_candidate(raw)
                if candidate and _is_valid_name(candidate):
                    holder_name = candidate.title()
                    break

        # ── Strategy 2: "Son / Daughter / Wife of" — the name is on the line before ──
        if not holder_name:
            so_pat = re.compile(
                r"([A-Z][A-Z\s\.]{1,60})\s*\n\s*(?:son|daughter|wife|husband|s/o|d/o|w/o)",
                re.IGNORECASE,
            )
            m = so_pat.search(cleaned)
            if m:
                candidate = _clean_candidate(m.group(1).strip())
                if candidate and _is_valid_name(candidate):
                    holder_name = candidate.title()

        # ── Strategy 3: scored generic scan ──────────────────────────────────
        if not holder_name:
            _name_candidates: list[tuple[str, int]] = []

            # Mixed-case pattern: "Ajai Kumar" / "N AJAI KUMAR" (allows single-letter initials)
            name_finder = re.compile(
                r"([A-Z](?:[a-zA-Z]{1,20})?(?:\s+[A-Z][a-zA-Z]{1,20})+)",
                re.MULTILINE,
            )
            for m in name_finder.finditer(cleaned):
                candidate = _clean_candidate(m.group(1).strip())
                words = candidate.split()
                if len(words) >= 2 and all(w.replace(".", "").isalpha() for w in words):
                    _name_candidates.append((candidate, m.start()))

            # UPPERCASE pattern: "AJAI KUMAR"
            uc_finder = re.compile(r"([A-Z]{1,20}(?:\s+[A-Z]{2,20})+)")
            for m in uc_finder.finditer(cleaned):
                candidate = _clean_candidate(m.group(1).strip())
                words = candidate.split()
                if len(words) >= 2 and all(len(w) >= 1 and w.isalpha() for w in words):
                    _name_candidates.append((candidate, m.start()))

            best_name = None
            best_score = -1
            for candidate, pos in _name_candidates:
                score = 0
                words = candidate.split()
                word_count = len(words)
                if 2 <= word_count <= 4:
                    score += 10
                elif word_count == 1:
                    continue
                if not _is_valid_name(candidate):
                    continue
                avg_len = sum(len(w) for w in words) / len(words)
                if avg_len < 2:
                    continue
                rel_pos = pos / max(len(cleaned), 1)
                if 0.3 < rel_pos < 0.8:
                    score += 5
                elif rel_pos >= 0.2:
                    score += 2
                after_text = cleaned[pos:pos + 100].lower()
                if re.search(r"son|daughter|wife|s/o|d/o|w/o", after_text):
                    score += 15
                before_text = cleaned[max(0, pos - 50):pos].lower()
                if re.search(r"\d{2}[-/]\d{2}[-/]\d{2,4}", before_text):
                    score += 8
                if score > best_score:
                    best_score = score
                    best_name = candidate

            if best_name:
                holder_name = best_name.title()

        # ── Extract blood group ──
        blood_group = None
        bg_match = re.search(r"(?:BLOOD\s*(?:GROUP)?|BG)\s*[:\-]?\s*((?:A|B|AB|O)[+-])", cleaned, re.IGNORECASE)
        if bg_match:
            blood_group = bg_match.group(1).upper()

        return {
            "license_number": license_number,
            "holder_name": holder_name,
            "issue_date": issue_date,
            "expiry_date": expiry_date,
            "dob": dob,
            "blood_group": blood_group,
        }

    def _normalize_license_token(self, value: str | None) -> str | None:
        if not value:
            return None
        compact = re.sub(r"[^A-Z0-9]", "", value.upper())
        return compact or None

    def _is_probable_driving_license_number(self, value: str | None) -> bool:
        compact = self._normalize_license_token(value)
        if not compact:
            return False
        if len(compact) < 10 or len(compact) > 18:
            return False
        if not re.match(r"^[A-Z]{2}\d", compact):
            return False
        if not re.search(r"\d{4}", compact):
            return False
        if compact in {"GOVERNMENT", "INDIA", "DRIVING", "LICENCE", "LICENSE"}:
            return False
        return True

    def _extract_driving_license_number_from_lines(self, lines: list[str]) -> str | None:
        line_pattern = re.compile(
            r"\b([A-Z]{2}\s*[-/]?\s*\d{1,2}\s*[-/]?\s*\d{4}\s*[-/]?\s*\d{4,10})\b",
            re.IGNORECASE,
        )

        for ln in lines[:12]:
            line_upper = ln.upper()

            # Direct match for common Indian DL pattern.
            for match in line_pattern.findall(line_upper):
                normalized = self._normalize_license_token(match)
                if self._is_probable_driving_license_number(normalized):
                    return normalized

            # OCR often splits DL as: "TN72 20240005499" or "TN72 2024 0005499".
            tokens = re.findall(r"[A-Z0-9]{2,}", line_upper)
            for i, token in enumerate(tokens):
                if not re.fullmatch(r"[A-Z]{2}\d{1,2}", token):
                    continue

                if i + 1 < len(tokens) and re.fullmatch(r"\d{9,12}", tokens[i + 1]):
                    candidate = self._normalize_license_token(token + tokens[i + 1])
                    if self._is_probable_driving_license_number(candidate):
                        return candidate

                if (
                    i + 2 < len(tokens)
                    and re.fullmatch(r"\d{4}", tokens[i + 1])
                    and re.fullmatch(r"\d{5,8}", tokens[i + 2])
                ):
                    candidate = self._normalize_license_token(token + tokens[i + 1] + tokens[i + 2])
                    if self._is_probable_driving_license_number(candidate):
                        return candidate

        # Try joining adjacent line pairs — OCR often splits "TN72" / "20240005499"
        # onto separate lines for smart card DLs.
        for i in range(min(len(lines) - 1, 11)):
            joined = lines[i].upper() + " " + lines[i + 1].upper()
            for match in line_pattern.findall(joined):
                normalized = self._normalize_license_token(match)
                if self._is_probable_driving_license_number(normalized):
                    return normalized
            # Also try 3-line join
            if i + 2 < len(lines):
                joined3 = joined + " " + lines[i + 2].upper()
                for match in line_pattern.findall(joined3):
                    normalized = self._normalize_license_token(match)
                    if self._is_probable_driving_license_number(normalized):
                        return normalized

        # Last fallback: generic alphanumeric token scan in upper section.
        generic_pattern = re.compile(r"\b[A-Z0-9\-/]{8,25}\b")
        for ln in lines[:12]:
            for token in generic_pattern.findall(ln.upper()):
                normalized = self._normalize_license_token(token)
                if self._is_probable_driving_license_number(normalized):
                    return normalized

        return None

    # ─────────────────────────────────────────────────────────
    # Passbook Tesseract extractor
    # ─────────────────────────────────────────────────────────

    def _extract_passbook_tesseract(
        self,
        file_bytes: bytes,
        media_type: str,
        document_type: str,
    ) -> dict:
        """Extract bank passbook fields (account number, IFSC, bank/branch name, holder name) via OCR."""
        try:
            from PIL import Image
            import io
            import pytesseract  # noqa: F401
        except Exception as e:
            return {"extracted": False, "reason": "tesseract_unavailable", "message": str(e)}

        if media_type == "application/pdf":
            try:
                file_bytes, _ = self._pdf_to_image(file_bytes)
            except Exception as e:
                return {"extracted": False, "reason": "pdf_conversion_error", "message": str(e)}

        try:
            image = self._preprocess_for_ocr(Image.open(io.BytesIO(file_bytes)))
            text = self._run_tesseract_multipass(image)
        except Exception as e:
            return {"extracted": False, "reason": "ocr_error", "message": str(e)}

        data = self._parse_passbook_text(text)
        if not any(v is not None for v in data.values() if isinstance(v, str)):
            return {
                "extracted": False,
                "reason": "no_fields",
                "message": "Could not detect bank passbook fields from the document.",
                "data": data,
            }
        return {"extracted": True, "document_type": document_type, "data": data}

    def _parse_passbook_text(self, text: str) -> dict:
        """Parse OCR text from an Indian bank passbook image."""
        lines = [ln.strip() for ln in (text or "").splitlines() if ln.strip()]
        joined = " ".join(lines)

        # ── Account number ──────────────────────────────────────────────────
        account_number = None
        acc_pat = re.compile(
            r"(?:account\s*(?:no|number|num)|a/?c\s*(?:no|number|num))[:\s.]+(\d[\d\s]{8,20}\d)",
            re.IGNORECASE,
        )
        m = acc_pat.search(joined)
        if m:
            account_number = re.sub(r"\s+", "", m.group(1))
        if not account_number:
            # Fallback: long digit sequence (10–18 digits) not a phone
            for tok in re.findall(r"\b\d{10,18}\b", joined):
                if not re.match(r"^[6-9]\d{9}$", tok):  # skip mobile numbers
                    account_number = tok
                    break

        # ── IFSC code ───────────────────────────────────────────────────────
        ifsc_code = None
        ifsc_pat = re.compile(r"\b([A-Z]{4}0[A-Z0-9]{6})\b", re.IGNORECASE)
        m = re.search(r"IFSC\s*[:\-]?\s*([A-Z0-9]{11})", joined, re.IGNORECASE)
        if m:
            candidate = m.group(1).upper()
            if re.match(r"^[A-Z]{4}0[A-Z0-9]{6}$", candidate):
                ifsc_code = candidate
        if not ifsc_code:
            m = ifsc_pat.search(joined)
            if m:
                ifsc_code = m.group(1).upper()

        # ── MICR code ───────────────────────────────────────────────────────
        micr_code = None
        m = re.search(r"MICR\s*[:\-]?\s*(\d{9})", joined, re.IGNORECASE)
        if m:
            micr_code = m.group(1)

        # ── Bank name ────────────────────────────────────────────────────────
        bank_name = None
        _KNOWN_BANKS = [
            "STATE BANK OF INDIA", "SBI", "HDFC BANK", "ICICI BANK",
            "AXIS BANK", "KOTAK MAHINDRA BANK", "PUNJAB NATIONAL BANK", "PNB",
            "BANK OF BARODA", "BANK OF INDIA", "CANARA BANK", "UNION BANK",
            "CENTRAL BANK OF INDIA", "INDIAN BANK", "INDIAN OVERSEAS BANK",
            "UCO BANK", "IDBI BANK", "YES BANK", "INDUSIND BANK",
            "FEDERAL BANK", "KARUR VYSYA BANK", "KVB", "SOUTH INDIAN BANK",
            "CITY UNION BANK", "CUB", "TAMILNAD MERCANTILE BANK", "TMB",
            "DHANLAXMI BANK", "LAKSHMI VILAS BANK", "KARNATAKA BANK",
            "DCB BANK", "RBL BANK", "BANDHAN BANK", "IDFC FIRST BANK",
            "AU SMALL FINANCE BANK",
        ]
        upper_joined = joined.upper()
        for bank in _KNOWN_BANKS:
            if bank in upper_joined:
                bank_name = bank
                break
        if not bank_name:
            # Use first content line that looks like a bank header
            for ln in lines[:5]:
                if re.search(r"BANK|FINANCE|CREDIT|COOPERATIVE", ln, re.IGNORECASE):
                    bank_name = ln.upper().strip()
                    break

        # ── Branch name ──────────────────────────────────────────────────────
        branch_name = None
        m = re.search(r"(?:branch\s*(?:name)?)[:\s]+([A-Za-z0-9\s,\.]{3,60}?)(?=\s*(?:IFSC|MICR|Account|\d{9,}|$))", joined, re.IGNORECASE)
        if m:
            branch_name = m.group(1).strip().upper()
            branch_name = re.sub(r"\s+", " ", branch_name).strip(",")

        # ── Account holder name ───────────────────────────────────────────────
        # Passbooks often put "Name:" on one line and the actual name on the next.
        account_holder_name = None

        # Strategy 1: same-line label (e.g. "Name of Account Holder: Sivabalan R")
        name_pat = re.compile(
            r"(?:(?:a/?c|account)\s*(?:holder|name)|name\s+of\s+account\s+holder|customer\s+name|depositor)[:\s]+([A-Za-z\s\.]{3,60})",
            re.IGNORECASE,
        )
        m = name_pat.search(joined)
        if m:
            raw = m.group(1).strip()
            raw = re.split(r"\d|Branch|Address|IFSC|Account|MICR|Customer|Code|Open|Date", raw, flags=re.IGNORECASE)[0].strip()
            raw = re.sub(r"\b(Mr|Mrs|Ms|Dr|Sri|Smt)\.?\s*", "", raw, flags=re.IGNORECASE).strip()
            # Validate: must be name-like (letters/spaces/dots/hyphens, ≥ 3 chars)
            if raw and re.match(r"^[A-Za-z][A-Za-z\s\.\-]*$", raw) and len(raw) >= 3:
                account_holder_name = raw.upper()

        # Strategy 2: next-line layout — "Name:" / "Name" label alone on a line,
        # name is on the very next non-empty line.
        if not account_holder_name:
            _name_label_pat = re.compile(
                r"^(?:name|a/?c\s*name|account\s*(?:holder|name)|customer\s*name|depositor)\s*:?\s*$",
                re.IGNORECASE,
            )
            for i, ln in enumerate(lines):
                if _name_label_pat.match(ln.strip()):
                    # Check next line(s)
                    for j in range(i + 1, min(i + 4, len(lines))):
                        candidate = lines[j].strip()
                        if not candidate:
                            continue
                        # Skip lines that look like labels (contain ':' or are known field names)
                        if ':' in candidate:
                            continue
                        if re.search(r"\b(IFSC|MICR|Account|Branch|Address|Customer|Open|Date|Code)\b", candidate, re.IGNORECASE):
                            continue
                        # Strip common honorifics
                        candidate = re.sub(r"\b(Mr|Mrs|Ms|Dr|Sri|Smt)\.?\s*", "", candidate, flags=re.IGNORECASE).strip()
                        # Must be letters/spaces/dots/hyphens only (name-like), ≥ 3 chars
                        if candidate and re.match(r"^[A-Za-z][A-Za-z\s\.\-]*$", candidate) and len(candidate) >= 3:
                            account_holder_name = candidate.upper()
                            break
                    if account_holder_name:
                        break

        # ── Bank address ─────────────────────────────────────────────────────
        bank_address = None
        m = re.search(r"(?:address)[:\s]+(.{5,120}?)(?:\n|$)", joined, re.IGNORECASE)
        if m:
            bank_address = re.sub(r"\s+", " ", m.group(1)).strip()

        # ── Confidence ───────────────────────────────────────────────────────
        filled = sum(1 for v in [account_holder_name, account_number, bank_name, branch_name, ifsc_code] if v)
        confidence = "high" if filled >= 4 else "medium" if filled >= 2 else "low"

        return {
            "account_holder_name": account_holder_name,
            "account_number": account_number,
            "bank_name": bank_name,
            "branch_name": branch_name,
            "ifsc_code": ifsc_code,
            "micr_code": micr_code,
            "bank_address": bank_address,
            "confidence": confidence,
        }

    def _extract_labeled_date(self, text: str, label_regex: str) -> str | None:
        pattern = re.compile(label_regex + r"\s*[:\-]?\s*([0-9A-Za-z\-/\. ]{6,20})", re.IGNORECASE)
        m = pattern.search(text)
        if not m:
            return None
        return self._normalize_date(m.group(1))

    def _extract_date_from_line_context(self, lines: list[str], label_patterns: list[str], pick: str = "first") -> str | None:
        compiled = [re.compile(p, re.IGNORECASE) for p in label_patterns]
        date_finder = re.compile(r"\b(?:\d{1,2}[\-/][A-Za-z]{3}[\-/]\d{2,4}|\d{2}[\-/]\d{2}[\-/]\d{4}|\d{4}[\-/]\d{2}[\-/]\d{2}|\d{2}\s+[A-Za-z]{3,9}\s+\d{4})\b")

        for i, line in enumerate(lines):
            if not any(p.search(line) for p in compiled):
                continue

            # Check current line first, then adjacent lines (OCR often wraps dates).
            window = [line]
            if i + 1 < len(lines):
                window.append(lines[i + 1])
            if i > 0:
                window.append(lines[i - 1])

            for w in window:
                found_dates: list[str] = []
                for token in date_finder.findall(w):
                    normalized = self._normalize_date(token)
                    if normalized:
                        found_dates.append(normalized)
                if found_dates:
                    ordered = self._sort_dates_unique(found_dates)
                    return ordered[-1] if pick == "latest" else ordered[0]
        return None

    def _extract_labeled_dates(self, text: str, label_regex: str) -> list[str]:
        pattern = re.compile(label_regex + r"\s*[:\-]?\s*([0-9A-Za-z\-/\. ]{6,20})", re.IGNORECASE)
        matches = pattern.findall(text)
        out: list[str] = []
        for m in matches:
            for token in re.findall(
                r"\b(?:\d{1,2}[\-/][A-Za-z]{3}[\-/]\d{2,4}|\d{2}[\-/]\d{2}[\-/]\d{4}|\d{4}[\-/]\d{2}[\-/]\d{2}|\d{2}\s+[A-Za-z]{3,9}\s+\d{4})\b",
                m,
            ):
                normalized = self._normalize_date(token)
                if normalized:
                    out.append(normalized)
        return out

    def _extract_all_dates(self, text: str) -> list[str]:
        candidates = re.findall(
            r"\b(?:\d{1,2}[\-/][A-Za-z]{3}[\-/]\d{2,4}|\d{2}[\-/]\d{2}[\-/]\d{4}|\d{4}[\-/]\d{2}[\-/]\d{2}|\d{2}\s+[A-Za-z]{3,9}\s+\d{4})\b",
            text,
        )
        out: list[str] = []
        for c in candidates:
            normalized = self._normalize_date(c)
            if normalized:
                out.append(normalized)
        return out

    def _normalize_date(self, raw: str | None) -> str | None:
        if not raw:
            return None
        value = re.sub(r"\s+", " ", raw.strip())
        value = value.replace(".", "-")
        value = value.replace("–", "-").replace("—", "-")
        value = re.sub(r"[^0-9A-Za-z\-/ ]", "", value)
        value = re.sub(r"\s*[-/]\s*", "-", value)
        value = re.sub(r"\s+", " ", value).strip()
        formats = [
            "%d/%m/%Y",
            "%d-%m-%Y",
            "%Y/%m/%d",
            "%Y-%m-%d",
            "%d-%b-%Y",
            "%d-%b-%y",
            "%d-%B-%Y",
            "%d-%B-%y",
            "%d %b %Y",
            "%d %B %Y",
            "%d %b %y",
            "%d %B %y",
            "%m-%d-%Y",
            "%m/%d/%Y",
        ]
        for fmt in formats:
            try:
                dt = datetime.strptime(value, fmt)
                return dt.strftime("%d/%m/%Y")
            except ValueError:
                continue
        return None

    def _sort_dates_unique(self, dates: list[str]) -> list[str]:
        parsed: list[tuple[datetime, str]] = []
        seen: set[str] = set()
        for d in dates:
            if d in seen:
                continue
            seen.add(d)
            try:
                parsed.append((datetime.strptime(d, "%d/%m/%Y"), d))
            except ValueError:
                continue
        parsed.sort(key=lambda x: x[0])
        return [d for _, d in parsed]

    def _pick_latest_date(self, dates: list[str]) -> str | None:
        ordered = self._sort_dates_unique(dates)
        if not ordered:
            return None
        return ordered[-1]

    def _pdf_to_image(self, pdf_bytes: bytes) -> tuple[bytes, str]:
        """Convert first page of a PDF to JPEG bytes."""
        from pdf2image import convert_from_bytes
        import io
        pages = convert_from_bytes(
            pdf_bytes, first_page=1, last_page=1, dpi=200
        )
        buf = io.BytesIO()
        pages[0].save(buf, format="JPEG", quality=90)
        return buf.getvalue(), "image/jpeg"

    def _pdf_to_images(self, pdf_bytes: bytes, max_pages: int = 3, dpi: int = 300) -> list[bytes]:
        """Convert first N pages of a PDF to JPEG bytes."""
        from pdf2image import convert_from_bytes
        import io

        pages = convert_from_bytes(
            pdf_bytes,
            first_page=1,
            last_page=max_pages,
            dpi=dpi,
        )
        if not pages:
            return []

        out: list[bytes] = []
        for page in pages:
            buf = io.BytesIO()
            page.save(buf, format="JPEG", quality=92)
            out.append(buf.getvalue())
        return out

    @staticmethod
    def _strip_code_fences(text: str) -> str:
        """Remove markdown ```json ... ``` wrapping if present."""
        text = text.strip()
        if text.startswith("```"):
            text = re.sub(r"^```(?:json)?\s*", "", text)
            text = re.sub(r"\s*```$", "", text)
        return text.strip()
